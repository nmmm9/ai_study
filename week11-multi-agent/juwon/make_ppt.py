from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BG     = RGBColor(0x0d, 0x11, 0x17)
BLUE   = RGBColor(0x58, 0xa6, 0xff)
GREEN  = RGBColor(0x3f, 0xb9, 0x50)
YELLOW = RGBColor(0xd2, 0x9e, 0x22)
RED    = RGBColor(0xf8, 0x51, 0x49)
GRAY   = RGBColor(0x8b, 0x94, 0x9e)
WHITE  = RGBColor(0xc9, 0xd1, 0xd9)
CARD   = RGBColor(0x16, 0x1b, 0x22)
DARK   = RGBColor(0x21, 0x26, 0x2d)

W, H = Inches(13.33), Inches(7.5)

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def bg(slide):
    s = slide.shapes.add_shape(1, 0, 0, W, H)
    s.fill.solid(); s.fill.fore_color.rgb = BG; s.line.fill.background()

def box(slide, x, y, w, h, color=CARD):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()

def line(slide, x, y, w, color=BLUE):
    s = slide.shapes.add_shape(1, x, y, w, Pt(3))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()

def txt(slide, text, x, y, w, h, size=14, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.color.rgb = color
    r.font.bold = bold;     r.font.name = "Malgun Gothic"


# ── 1. 표지 ──────────────────────────────────────────────────
def slide_cover(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    line(sl, Inches(1), Inches(3.2), Inches(11.3), BLUE)
    txt(sl, "Week 11", Inches(1), Inches(1.2), Inches(11), Inches(1),
        size=44, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    txt(sl, "Multi-Agent Debate", Inches(1), Inches(2.2), Inches(11), Inches(0.85),
        size=28, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(sl, "GitHub Tech Trend Analyzer — 여러 AI 전문가가 토론해서 트렌드를 분석하는 시스템",
        Inches(1), Inches(3.35), Inches(11), Inches(0.6),
        size=14, color=GRAY, align=PP_ALIGN.CENTER)
    tags = [
        ("에이전트 구조", BLUE), ("각 에이전트 역할", GREEN),
        ("Debate 패턴이란?", YELLOW), ("10주차 vs 11주차", RED), ("오늘 추가 기능", GREEN)
    ]
    for i, (tag, c) in enumerate(tags):
        box(sl, Inches(0.8 + i * 2.38), Inches(4.3), Inches(2.2), Inches(0.6), DARK)
        txt(sl, tag, Inches(0.85 + i * 2.38), Inches(4.4), Inches(2.1), Inches(0.45),
            size=12, color=c, bold=True, align=PP_ALIGN.CENTER)


# ── 2. 멀티 에이전트 시스템이란? ─────────────────────────────
def slide_mas_intro(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    txt(sl, "멀티 에이전트 시스템이란?", Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
        size=26, color=BLUE, bold=True)
    line(sl, Inches(0.5), Inches(0.8), Inches(12.3), BLUE)

    box(sl, Inches(0.3), Inches(1.0), Inches(12.5), Inches(0.85), CARD)
    txt(sl, "여러 AI 에이전트가 각자의 역할을 맡아 협력하면서 하나의 목표를 달성하는 시스템",
        Inches(0.5), Inches(1.12), Inches(12.1), Inches(0.55), size=15, color=WHITE, bold=True)

    # 3대 구성요소
    components = [
        ("🤖 에이전트", BLUE,
         "시스템 내에서 독립적으로 행동하는 개체\n\n• 자율성을 갖고 환경을 인식하여 판단\n• 소프트웨어 프로그램, 봇, 로봇 등\n\n우리 프로젝트:\nagent_ai / agent_web / agent_sec /\nSupervisor / Critic / Judge"),
        ("🌍 환경", GREEN,
         "에이전트들이 공유하며 일하는 공간\n\n• 리소스를 제공하고 제약 조건 적용\n• 에이전트 간 간접 소통의 매개체\n\n우리 프로젝트:\nLangGraph StateGraph\n(모든 에이전트가 공유하는 상태 객체)"),
        ("📡 통신", YELLOW,
         "에이전트끼리 정보를 주고받는 규칙\n\n• 메시지 형식과 전달 방식 정의\n• FIPA ACL 같은 표준 언어 존재\n\n우리 프로젝트:\n각 노드가 State를 읽고 결과를 State에\n저장 → 다음 노드가 이를 받아 사용"),
    ]
    for i, (title, c, desc) in enumerate(components):
        bx = Inches(0.3 + i * 4.35)
        box(sl, bx, Inches(2.0), Inches(4.1), Inches(5.1), CARD)
        txt(sl, title, bx + Inches(0.2), Inches(2.15), Inches(3.7), Inches(0.5), size=16, color=c, bold=True)
        line(sl, bx + Inches(0.2), Inches(2.7), Inches(3.7), DARK)
        txt(sl, desc, bx + Inches(0.2), Inches(2.8), Inches(3.7), Inches(4.0), size=10.5, color=WHITE)

    box(sl, Inches(0.3), Inches(7.1), Inches(12.5), Inches(0.22), DARK)
    txt(sl, "핵심: 한 명의 AI가 모든 걸 판단하는 게 아니라, 역할을 나눈 여러 AI가 협력해서 더 좋은 결론을 냅니다",
        Inches(0.5), Inches(7.12), Inches(12.1), Inches(0.2), size=11, color=GRAY)


# ── 3. MAS 아키텍처 & 구조 ────────────────────────────────────
def slide_mas_arch(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    txt(sl, "MAS 아키텍처 & 구조 유형", Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
        size=26, color=GREEN, bold=True)
    line(sl, Inches(0.5), Inches(0.8), Inches(12.3), GREEN)

    # 네트워크 구조 (왼쪽)
    txt(sl, "네트워크 구조 — 어떻게 연결되나?", Inches(0.5), Inches(1.0), Inches(6), Inches(0.4),
        size=13, color=YELLOW, bold=True)
    nets = [
        ("중앙 집중식", BLUE,
         "중앙 컨트롤타워가 모든 에이전트를 관리\n장점: 소통 쉽고 지식 통일\n단점: 중앙 서버 죽으면 전체 다운"),
        ("분산형", GRAY,
         "에이전트가 이웃끼리 직접 소통\n장점: 한 에이전트 죽어도 시스템 유지\n단점: 전체 조율이 어려움"),
    ]
    for i, (title, c, desc) in enumerate(nets):
        bx = Inches(0.3 + i * 3.1)
        box(sl, bx, Inches(1.5), Inches(2.9), Inches(2.8), CARD)
        txt(sl, title, bx + Inches(0.15), Inches(1.65), Inches(2.6), Inches(0.4), size=12, color=c, bold=True)
        line(sl, bx + Inches(0.15), Inches(2.1), Inches(2.6), DARK)
        txt(sl, desc, bx + Inches(0.15), Inches(2.2), Inches(2.6), Inches(1.1), size=10, color=WHITE)

    # 조직 구조 (오른쪽)
    txt(sl, "조직 구조 — 어떻게 역할을 나누나?", Inches(6.5), Inches(1.0), Inches(6.3), Inches(0.4),
        size=13, color=YELLOW, bold=True)
    structs = [
        ("계층적 구조", BLUE,   "상사→부하 처럼 위에서 아래로\n권한과 책임이 명확하게 분리됨"),
        ("Holonic",   GREEN,  "팀 안에 또 팀이 있는 구조\n하위 에이전트가 여러 팀에 속할 수 있음"),
        ("연합",       YELLOW, "필요할 때만 팀 꾸리고 해산\n유연하지만 조율이 어려움"),
        ("팀",         GRAY,   "한 목표를 위해 서로 의존하며 협력\n연합보다 더 긴밀한 관계"),
    ]
    for i, (title, c, desc) in enumerate(structs):
        col = i % 2
        row = i // 2
        bx = Inches(6.5 + col * 3.2)
        by = Inches(1.5 + row * 1.7)
        box(sl, bx, by, Inches(3.0), Inches(1.5), CARD)
        txt(sl, title, bx + Inches(0.15), by + Inches(0.1), Inches(2.7), Inches(0.4), size=12, color=c, bold=True)
        txt(sl, desc,  bx + Inches(0.15), by + Inches(0.55), Inches(2.7), Inches(0.85), size=10, color=WHITE)

    # 우리 프로젝트 적용
    box(sl, Inches(0.3), Inches(5.1), Inches(12.5), Inches(1.75), DARK)
    txt(sl, "우리 프로젝트가 사용한 구조", Inches(0.5), Inches(5.2), Inches(12), Inches(0.4),
        size=13, color=GREEN, bold=True)
    line(sl, Inches(0.5), Inches(5.65), Inches(12.1), DARK)
    txt(sl, "✅  중앙 집중식 + 계층적 구조\n"
            "LangGraph가 중앙에서 전체 흐름을 관리하고, collect → 전문가 3명 → Supervisor → Critic → Judge 순으로\n"
            "위에서 아래로 권한이 내려가는 계층적 구조로 설계했습니다.",
        Inches(0.5), Inches(5.72), Inches(12.1), Inches(1.0), size=11, color=WHITE)


# ── 4. 에이전트 구조도 (구 2번) ───────────────────────────────
def slide_structure(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    txt(sl, "에이전트 구조도", Inches(0.5), Inches(0.2), Inches(7), Inches(0.6),
        size=26, color=BLUE, bold=True)
    txt(sl, "LangGraph의 StateGraph로 구현 — 각 노드가 하나의 에이전트",
        Inches(0.5), Inches(0.8), Inches(10), Inches(0.45), size=13, color=GRAY)
    line(sl, Inches(0.5), Inches(1.2), Inches(12.3), BLUE)

    # 노드들
    nodes = [
        (Inches(0.3),  Inches(1.5), "📡 collect",    "GitHub Trending\n페이지에서 레포 수집",        GRAY),
        (Inches(3.0),  Inches(1.0), "🧠 agent_ai",   "AI/ML 관련 레포\n전문가 시각으로 분석",        BLUE),
        (Inches(3.0),  Inches(2.5), "🌐 agent_web",  "웹/앱 관련 레포\n전문가 시각으로 분석",        GREEN),
        (Inches(3.0),  Inches(4.0), "🔒 agent_sec",  "보안 관련 레포\n전문가 시각으로 분석",        YELLOW),
        (Inches(6.0),  Inches(2.5), "📊 supervisor", "3개 분석을\n하나로 종합",                     BLUE),
        (Inches(8.7),  Inches(2.5), "⚡ critic",     "종합 결과에\n반론 제기",                       RED),
        (Inches(6.0),  Inches(4.8), "⚖️ judge",      "Supervisor + Critic\n검토 후 최종 결론",      GREEN),
        (Inches(9.5),  Inches(4.8), "💾 report",     "결과 저장\n메일 + GitHub 업로드",             GRAY),
    ]
    for (x, y, name, desc, c) in nodes:
        box(sl, x, y, Inches(2.5), Inches(1.1), CARD)
        txt(sl, name, x + Inches(0.15), y + Inches(0.1), Inches(2.2), Inches(0.45), size=12, color=c, bold=True)
        txt(sl, desc, x + Inches(0.15), y + Inches(0.55), Inches(2.2), Inches(0.5), size=10, color=WHITE)

    # 화살표 표시
    arrows = [
        (Inches(2.85), Inches(1.85), "→"),
        (Inches(2.85), Inches(3.35), "→"),
        (Inches(2.85), Inches(4.85), "→"),
        (Inches(5.55), Inches(3.05), "→"),
        (Inches(8.25), Inches(3.05), "→"),
        (Inches(7.25), Inches(4.55), "↓"),
        (Inches(11.3), Inches(5.1), "→"),  # judge → report
    ]
    for (x, y, arrow) in arrows:
        txt(sl, arrow, x, y, Inches(0.45), Inches(0.4), size=16, color=GRAY, align=PP_ALIGN.CENTER)

    # critic에서 judge로
    txt(sl, "↓", Inches(9.75), Inches(4.2), Inches(0.45), Inches(0.45), size=16, color=GRAY, align=PP_ALIGN.CENTER)

    box(sl, Inches(0.3), Inches(6.3), Inches(12.5), Inches(0.6), DARK)
    txt(sl, "💡  agent_ai / agent_web / agent_sec 는 동시에(병렬로) 실행됩니다 — LangGraph의 fan-out 기능 활용",
        Inches(0.5), Inches(6.38), Inches(12), Inches(0.45), size=12, color=YELLOW)


# ── 3. 각 에이전트 역할 ───────────────────────────────────────
def slide_roles(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    txt(sl, "각 에이전트 역할", Inches(0.5), Inches(0.2), Inches(10), Inches(0.6),
        size=26, color=GREEN, bold=True)
    line(sl, Inches(0.5), Inches(0.8), Inches(12.3), GREEN)

    agents = [
        ("🧠 agent_ai",   "AI/ML 전문가",  BLUE,
         "키워드 필터: ai, ml, llm, gpt, neural, deep learning\n→ AI/ML 관련 레포만 추려서 기술 트렌드 분석\n→ 급부상 중인 AI 기술 및 추천 레포 2개 선정"),
        ("🌐 agent_web",  "웹/앱 전문가",  GREEN,
         "키워드 필터: react, vue, next, typescript, frontend, mobile\n→ 웹/앱 개발 관련 레포만 추려서 분석\n→ 주목할 프레임워크·라이브러리 및 추천 레포 2개 선정"),
        ("🔒 agent_sec",  "보안 전문가",   YELLOW,
         "키워드 필터: security, auth, crypto, vulnerability, pentest\n→ 보안 관련 레포만 추려서 분석\n→ 급부상 중인 보안 이슈 및 추천 레포 2개 선정"),
        ("📊 supervisor", "종합 분석가",   BLUE,
         "agent_ai + agent_web + agent_sec 3개 분석을 모두 입력으로 받음\n→ 세 분야를 아우르는 종합 트렌드 리포트 작성\n→ 전체적인 기술 흐름과 핵심 인사이트 정리"),
        ("⚡ critic",     "비판적 검토자", RED,
         "supervisor 리포트를 입력으로 받음\n→ 빠진 트렌드, 근거 부족한 주장, 더 강조할 부분 지적\n→ 반론과 보완점을 구체적으로 제시"),
        ("⚖️ judge",      "최종 심판자",   GREEN,
         "supervisor 리포트 + critic 반론을 모두 입력으로 받음\n→ 두 의견을 종합해서 최종 트렌드 리포트 완성\n→ 편향 없이 균형 잡힌 결론 도출"),
    ]

    for i, (name, role, c, desc) in enumerate(agents):
        col = i % 3
        row = i // 3
        bx = Inches(0.3 + col * 4.35)
        by = Inches(1.0 + row * 3.0)
        box(sl, bx, by, Inches(4.1), Inches(2.7), CARD)
        txt(sl, name, bx + Inches(0.2), by + Inches(0.15), Inches(3.7), Inches(0.45), size=13, color=c, bold=True)
        txt(sl, role, bx + Inches(0.2), by + Inches(0.6),  Inches(3.7), Inches(0.3),  size=11, color=GRAY)
        line(sl, bx + Inches(0.2), by + Inches(0.95), Inches(3.7), DARK)
        txt(sl, desc, bx + Inches(0.2), by + Inches(1.05), Inches(3.7), Inches(1.5),  size=10, color=WHITE)


# ── 4. 실제 에이전트 출력 예시 ────────────────────────────────
def slide_example(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    txt(sl, "실제 에이전트 출력 예시", Inches(0.5), Inches(0.2), Inches(10), Inches(0.6),
        size=26, color=YELLOW, bold=True)
    txt(sl, "Python 언어, 이번 주 기준으로 분석했을 때 실제로 나오는 내용",
        Inches(0.5), Inches(0.75), Inches(12), Inches(0.4), size=13, color=GRAY)
    line(sl, Inches(0.5), Inches(1.1), Inches(12.3), YELLOW)

    examples = [
        ("🧠 agent_ai 출력", BLUE,
         "주목할 AI/ML 트렌드:\n"
         "• TauricResearch/TradingAgents — LLM 기반 금융 트레이딩 프레임워크로\n"
         "  멀티 에이전트 협업을 금융 도메인에 적용한 점이 주목됨\n"
         "• anthropics/financial-services — 금융 서비스용 Claude 플러그인.\n"
         "  AI가 실제 산업에 빠르게 통합되는 트렌드를 보여줌"),
        ("⚡ critic 반론", RED,
         "Supervisor 리포트 검토 결과:\n"
         "• 보안 측면이 과소 평가됨 — AI 금융 시스템의 취약점과\n"
         "  데이터 프라이버시 이슈가 충분히 다뤄지지 않음\n"
         "• 중국어 레포(TradingAgents-CN)의 비중이 크나\n"
         "  글로벌 트렌드로 일반화하기엔 근거 부족"),
        ("⚖️ judge 최종 결론", GREEN,
         "Critic의 보안 지적을 반영한 최종 결론:\n"
         "• AI × 금융 분야가 이번 주 가장 뜨거운 트렌드\n"
         "• TradingAgents, financial-services-plugins 모두 주목 필요\n"
         "• 단, 금융 AI 도입 시 보안·규제 준수 검토가 선행되어야 함"),
    ]

    for i, (title, c, content) in enumerate(examples):
        bx = Inches(0.3 + i * 4.35)
        box(sl, bx, Inches(1.5), Inches(4.1), Inches(4.85), CARD)
        txt(sl, title, bx + Inches(0.2), Inches(1.65), Inches(4.5), Inches(0.45), size=12, color=c, bold=True)
        line(sl, bx + Inches(0.2), Inches(2.1), Inches(4.5), DARK)
        txt(sl, content, bx + Inches(0.2), Inches(2.2), Inches(4.5), Inches(3.6), size=10, color=WHITE)

    box(sl, Inches(0.3), Inches(6.55), Inches(12.5), Inches(0.55), DARK)
    txt(sl, "💬 채팅 질문 예시:  \"이번주 가장 주목할 레포 하나만 뽑아줘\"  →  AI가 위 분석 결과를 바탕으로 답변",
        Inches(0.5), Inches(6.62), Inches(12), Inches(0.42), size=11, color=GRAY)


# ── 5. Debate 패턴이란? ───────────────────────────────────────
def slide_debate(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    txt(sl, "Debate 패턴이란?", Inches(0.5), Inches(0.2), Inches(10), Inches(0.6),
        size=26, color=YELLOW, bold=True)
    line(sl, Inches(0.5), Inches(0.8), Inches(12.3), YELLOW)

    # 개념 설명
    box(sl, Inches(0.3), Inches(1.0), Inches(12.5), Inches(1.3), CARD)
    txt(sl, "멀티 에이전트 패턴 중 하나로, 여러 AI가 서로 다른 역할을 맡아 '토론'하면서 결론을 도출하는 방식입니다.",
        Inches(0.5), Inches(1.1), Inches(12.1), Inches(0.45), size=14, color=WHITE, bold=True)
    txt(sl, "한 명의 AI가 혼자 판단하면 편향이 생길 수 있어서, 주장→반론→최종 판단 구조로 서로 검증합니다.",
        Inches(0.5), Inches(1.6), Inches(12.1), Inches(0.4), size=13, color=GRAY)

    # 3단계
    steps = [
        ("1단계\n주장", "Supervisor가 3명의 전문가 분석을 종합해서\n리포트를 작성합니다.\n\n\"이번 주 AI 금융 트렌드가 가장 주목됩니다\"", BLUE),
        ("2단계\n반론", "Critic이 리포트를 검토하고\n부족한 부분을 지적합니다.\n\n\"보안 이슈가 빠졌고 근거가 부족합니다\"", RED),
        ("3단계\n최종 판단", "Judge가 주장과 반론을 모두 보고\n균형 잡힌 최종 결론을 냅니다.\n\n\"보안 주의사항을 포함한 최종 리포트\"", GREEN),
    ]
    for i, (title, desc, c) in enumerate(steps):
        bx = Inches(0.3 + i * 4.35)
        box(sl, bx, Inches(2.5), Inches(4.0), Inches(3.8), CARD)
        txt(sl, title, bx + Inches(0.2), Inches(2.65), Inches(3.4), Inches(0.7), size=16, color=c, bold=True)
        line(sl, bx + Inches(0.2), Inches(3.4), Inches(3.4), DARK)
        txt(sl, desc,  bx + Inches(0.2), Inches(3.5), Inches(3.4), Inches(2.8), size=11, color=WHITE)
        if i < 2:
            txt(sl, "→", Inches(4.15 + i * 4.35), Inches(3.8), Inches(0.45), Inches(0.5),
                size=20, color=GRAY, align=PP_ALIGN.CENTER)

    # 다른 패턴 비교
    box(sl, Inches(0.3), Inches(6.55), Inches(12.5), Inches(0.65), DARK)
    patterns = "Pipeline(순차전달)  |  Supervisor(지휘)  |  ✅ Debate(토론) — 이번 주  |  Voting(다수결)  |  ReAct(행동+관찰)"
    txt(sl, "다른 패턴들: " + patterns, Inches(0.5), Inches(6.62), Inches(12), Inches(0.5), size=11, color=GRAY)


# ── 6. 10주차 vs 11주차 ───────────────────────────────────────
def slide_vs(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    txt(sl, "10주차 vs 11주차", Inches(0.5), Inches(0.2), Inches(10), Inches(0.6),
        size=26, color=RED, bold=True)
    line(sl, Inches(0.5), Inches(0.8), Inches(12.3), RED)

    # 10주차
    box(sl, Inches(0.3), Inches(1.0), Inches(5.9), Inches(6.1), CARD)
    txt(sl, "10주차 — Self-Correction", Inches(0.5), Inches(1.15), Inches(5.5), Inches(0.5), size=17, color=BLUE, bold=True)
    txt(sl, "핵심 개념: AI가 스스로 결과물의 품질을 채점하고\n기준 미달이면 피드백을 반영해 다시 작성",
        Inches(0.5), Inches(1.7), Inches(5.5), Inches(0.7), size=11, color=GRAY)
    line(sl, Inches(0.5), Inches(2.45), Inches(5.5), DARK)
    w10 = [
        ("수집", "GitHub에서 트렌딩 레포 수집", WHITE),
        ("생성", "AI가 트렌드 분석 보고서 작성", WHITE),
        ("채점", "AI가 직접 0~100점으로 품질 평가\n(5가지 기준: 구체성, 인사이트, 구조 등)", YELLOW),
        ("판단", "70점 이상 → 통과\n70점 미만 → 피드백과 함께 재작성 요청", GREEN),
        ("완성", "최종 보고서 저장 + 메일 + GitHub 업로드", WHITE),
    ]
    for i, (step, desc, c) in enumerate(w10):
        txt(sl, f"{'①②③④⑤'[i]} {step}",  Inches(0.5), Inches(2.6 + i * 0.82), Inches(1.5), Inches(0.4), size=12, color=c, bold=True)
        txt(sl, desc, Inches(1.9), Inches(2.6 + i * 0.82), Inches(4.1), Inches(0.6), size=10, color=WHITE)

    txt(sl, "VS", Inches(6.25), Inches(3.6), Inches(0.8), Inches(0.7),
        size=24, color=RED, bold=True, align=PP_ALIGN.CENTER)

    # 11주차
    box(sl, Inches(7.1), Inches(1.0), Inches(5.9), Inches(6.1), CARD)
    txt(sl, "11주차 — Multi-Agent Debate", Inches(7.3), Inches(1.15), Inches(5.5), Inches(0.5), size=17, color=GREEN, bold=True)
    txt(sl, "핵심 개념: 여러 전문가 AI가 각자 분석하고 토론해서\n편향 없는 균형 잡힌 결론을 도출",
        Inches(7.3), Inches(1.7), Inches(5.5), Inches(0.7), size=11, color=GRAY)
    line(sl, Inches(7.3), Inches(2.45), Inches(5.5), DARK)
    w11 = [
        ("수집",    "GitHub Trending 스크래핑으로 실제 트렌드 수집",                        WHITE),
        ("병렬분석", "AI/ML·웹/앱·보안 전문가 3명이 동시에 각자 분야 분석",                   BLUE),
        ("종합",    "Supervisor가 3개 분석을 하나의 리포트로 합침",                          WHITE),
        ("반론",    "Critic이 리포트의 허점·편향·누락을 지적",                               RED),
        ("최종",    "Judge가 주장+반론 검토 후 균형 잡힌 최종 리포트 완성",                   GREEN),
    ]
    for i, (step, desc, c) in enumerate(w11):
        txt(sl, f"{'①②③④⑤'[i]} {step}", Inches(7.3), Inches(2.6 + i * 0.82), Inches(1.5), Inches(0.4), size=12, color=c, bold=True)
        txt(sl, desc, Inches(8.7), Inches(2.6 + i * 0.82), Inches(4.1), Inches(0.6), size=10, color=WHITE)


# ── 7. 오늘 추가된 기능 ───────────────────────────────────────
def slide_today(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    txt(sl, "오늘 추가된 기능", Inches(0.5), Inches(0.2), Inches(10), Inches(0.6),
        size=26, color=GREEN, bold=True)
    line(sl, Inches(0.5), Inches(0.8), Inches(12.3), GREEN)

    features = [
        ("①", "트렌드 변화 감지", GREEN,
         "이전 분석 결과와 현재를 자동 비교\n→ 🆕 새로 등장 / 🚀 급등 / 📉 급락 / 👋 사라짐 4가지로 분류\n→ 분석을 두 번 이상 실행하면 자동으로 표시"),
        ("②", "자동 스케줄링", BLUE,
         "APScheduler 라이브러리로 백그라운드에서 실행\n→ 매일 지정한 시간에 자동으로 분석 + 메일 발송\n→ 서버가 켜져 있는 동안 계속 동작"),
        ("③", "채팅 기능", YELLOW,
         "분석 결과 전체를 AI에게 컨텍스트로 전달\n→ 결과를 보면서 바로 질문 가능\n→ 예: \"보안 관련만 정리해줘\", \"쉽게 설명해줘\""),
        ("④", "UI 개선", RED,
         "react-markdown 라이브러리 도입\n→ ###, **굵게**, - 목록 등이 실제 스타일로 렌더링\n→ 전문가 분석·토론·Judge 결론 모두 적용"),
        ("⑤", "데이터 수집 개선", GREEN,
         "GitHub API 검색 → GitHub Trending 페이지 스크래핑으로 교체\n→ '이번 주 획득 스타 수' 직접 수집 가능\n→ OSSInsight와 유사한 정확도 달성"),
    ]

    for i, (num, title, c, desc) in enumerate(features):
        by = Inches(1.0 + i * 1.22)
        box(sl, Inches(0.3), by, Inches(12.5), Inches(1.1), CARD)
        txt(sl, num,   Inches(0.5),  by + Inches(0.2), Inches(0.5),  Inches(0.7), size=20, color=c, bold=True)
        txt(sl, title, Inches(1.1),  by + Inches(0.22), Inches(2.8), Inches(0.6), size=14, color=WHITE, bold=True)
        txt(sl, desc,  Inches(4.1),  by + Inches(0.12), Inches(8.5), Inches(0.88), size=11, color=GRAY)


# ── 8. 마무리 ─────────────────────────────────────────────────
def slide_end(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl)
    line(sl, Inches(1), Inches(3.4), Inches(11.3), BLUE)
    txt(sl, "혼자 판단하는 AI  →  토론하는 AI",
        Inches(1), Inches(1.5), Inches(11), Inches(0.9),
        size=30, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(sl, "전문화된 에이전트 3명이 각자 분석하고, Critic이 검증하고, Judge가 최종 결론을 내립니다.",
        Inches(1), Inches(2.5), Inches(11), Inches(0.6),
        size=15, color=GRAY, align=PP_ALIGN.CENTER)
    txt(sl, "더 넓은 시각  ·  서로 검증  ·  균형 잡힌 결론",
        Inches(1), Inches(3.6), Inches(11), Inches(0.55),
        size=16, color=BLUE, bold=True, align=PP_ALIGN.CENTER)

    keywords = [("전문화", BLUE), ("병렬 처리", GREEN), ("검증", YELLOW), ("균형", RED)]
    for i, (kw, c) in enumerate(keywords):
        box(sl, Inches(1.8 + i * 2.45), Inches(4.5), Inches(2.2), Inches(0.75), DARK)
        txt(sl, kw, Inches(1.85 + i * 2.45), Inches(4.6), Inches(2.1), Inches(0.55),
            size=17, color=c, bold=True, align=PP_ALIGN.CENTER)


# ── 실행 ─────────────────────────────────────────────────────
prs = new_prs()
slide_cover(prs)
slide_mas_intro(prs)
slide_mas_arch(prs)
slide_structure(prs)
slide_roles(prs)
slide_example(prs)
slide_debate(prs)
slide_vs(prs)
slide_today(prs)
slide_end(prs)

out = r"c:\윤주원\ai study\ai_study\week11-multi-agent\juwon\week11_발표_v3.pptx"
prs.save(out)
print("저장 완료:", out)
