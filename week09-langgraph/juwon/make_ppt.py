"""
9주차 발표용 PPT 생성
실행: python make_ppt.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BLUE     = RGBColor(0x26, 0x5D, 0xC8)
DARK     = RGBColor(0x1E, 0x1E, 0x2E)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GRAY     = RGBColor(0x94, 0x94, 0xB0)
LIGHT_BG = RGBColor(0xF4, 0xF6, 0xFF)
GREEN    = RGBColor(0x16, 0xA3, 0x4A)
ORANGE   = RGBColor(0xEA, 0x58, 0x0C)
PURPLE   = RGBColor(0x9B, 0x59, 0xB6)

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def txt(slide, text, left, top, width, height,
        size=16, bold=False, color=DARK, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


# ══════════════════════════════════════════════
# 슬라이드 1 — 표지
# ══════════════════════════════════════════════
s1 = prs.slides.add_slide(blank)
set_bg(s1, DARK)

box(s1, 0, 2.8, 0.12, 1.8, BLUE)
box(s1, 0.4, 1.4, 1.8, 0.45, BLUE)
txt(s1, "WEEK 09", 0.4, 1.4, 1.8, 0.45, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txt(s1, "LangGraph", 0.4, 2.0, 10, 0.9, size=52, bold=True, color=WHITE)
txt(s1, "입문", 0.4, 2.9, 12, 0.9, size=52, bold=True, color=BLUE)
txt(s1, "그래프 기반 워크플로우 (Nodes, Edges)", 0.4, 3.9, 10, 0.6, size=18, color=GRAY)
txt(s1, "GitHub 기술 트렌드 분석 에이전트  —  5개 노드 자동 실행", 0.4, 4.6, 10, 0.5, size=15, bold=True, color=BLUE)
txt(s1, "juwon  ·  2026.04.28", 0.4, 6.6, 6, 0.4, size=13, color=GRAY)


# ══════════════════════════════════════════════
# 슬라이드 2 — 8주차 vs 9주차
# ══════════════════════════════════════════════
s2 = prs.slides.add_slide(blank)
set_bg(s2, LIGHT_BG)
box(s2, 0, 0, 13.33, 1.3, BLUE)
txt(s2, "8주차 vs 9주차: 무엇이 달라졌나", 0.5, 0.2, 12, 1.0, size=30, bold=True, color=WHITE)

# 8주차
box(s2, 0.4, 1.45, 5.8, 5.1, RGBColor(0xFF,0xE5,0xE5))
txt(s2, "8주차  ReAct (while 루프)", 0.6, 1.6, 5.4, 0.5, size=16, bold=True, color=RGBColor(0xCC,0,0))
code8 = """while iteration < 25:
    response = client.chat(
        messages, tools=TOOLS
    )
    if tool_calls:
        execute_tool()
        messages.append(result)
    else:
        break"""
box(s2, 0.6, 2.2, 5.4, 2.8, DARK)
txt(s2, code8, 0.75, 2.3, 5.1, 2.6, size=11, color=RGBColor(0xA8,0xE6,0xFF))
items8 = ["흐름을 코드로만 파악 가능", "조건 분기가 복잡해짐", "시각화 불가"]
for i, t in enumerate(items8):
    txt(s2, f"✗  {t}", 0.6, 5.15 + i*0.38, 5.4, 0.35, size=12,
        color=RGBColor(0xCC,0,0))

# 9주차
box(s2, 7.0, 1.45, 5.9, 5.1, RGBColor(0xE5,0xFF,0xED))
txt(s2, "9주차  LangGraph (그래프)", 7.2, 1.6, 5.5, 0.5, size=16, bold=True, color=GREEN)
code9 = """graph.add_node("collect",  collect_node)
graph.add_node("validate", validate_node)
graph.add_node("analyze",  analyze_node)
graph.add_conditional_edges(
    "validate", should_retry,
    {"retry": "collect",
     "analyze": "analyze"}
)"""
box(s2, 7.2, 2.2, 5.5, 2.8, DARK)
txt(s2, code9, 7.35, 2.3, 5.2, 2.6, size=11, color=RGBColor(0xA8,0xE6,0xFF))
items9 = ["흐름을 그래프로 한눈에", "조건 분기가 한 줄로 표현", "시각화 자동 제공"]
for i, t in enumerate(items9):
    txt(s2, f"✓  {t}", 7.2, 5.15 + i*0.38, 5.5, 0.35, size=12, color=GREEN)

txt(s2, "핵심: while 루프를 노드와 엣지로 분리 → 흐름이 명확해진다",
    0.4, 6.75, 12.5, 0.4, size=14, bold=True, color=BLUE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════
# 슬라이드 3 — 왜 LangGraph인가? (장점)
# ══════════════════════════════════════════════
s3 = prs.slides.add_slide(blank)
set_bg(s3, LIGHT_BG)
box(s3, 0, 0, 13.33, 1.3, BLUE)
txt(s3, "어차피 실행은 같은데 왜 LangGraph를 쓰나?", 0.5, 0.2, 12, 1.0, size=28, bold=True, color=WHITE)

# 핵심 전제
box(s3, 0.4, 1.4, 12.5, 0.6, DARK)
txt(s3, "맞아요, 실행 결과는 똑같아요.  차이는 \"관리 방식\" 이에요.",
    0.6, 1.48, 12.1, 0.45, size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# 장점 1
box(s3, 0.4, 2.15, 3.9, 4.8, RGBColor(0xE8,0xF0,0xFF))
box(s3, 0.4, 2.15, 3.9, 0.5, BLUE)
txt(s3, "① 흐름이 한눈에 보임", 0.55, 2.18, 3.6, 0.42, size=13, bold=True, color=WHITE)
txt(s3, "while 루프:", 0.55, 2.78, 3.6, 0.35, size=11, bold=True, color=RGBColor(0xCC,0,0))
txt(s3, "코드 100줄 읽어야\n흐름 파악 가능", 0.55, 3.15, 3.6, 0.65, size=11, color=DARK)
txt(s3, "LangGraph:", 0.55, 3.9, 3.6, 0.35, size=11, bold=True, color=GREEN)
txt(s3, "그래프 그림 1장으로\n바로 파악", 0.55, 4.28, 3.6, 0.65, size=11, color=DARK)

# 장점 2
box(s3, 4.5, 2.15, 4.3, 4.8, RGBColor(0xF0,0xFF,0xF4))
box(s3, 4.5, 2.15, 4.3, 0.5, GREEN)
txt(s3, "② 단계 수정이 쉬움", 4.65, 2.18, 4.0, 0.42, size=13, bold=True, color=WHITE)
txt(s3, "\"분석 전에 번역 추가하고 싶어\"", 4.65, 2.78, 4.0, 0.42, size=11, color=DARK)
txt(s3, "while 루프:\n코드 전체 읽고\n어디에 넣을지 찾아야 함", 4.65, 3.28, 4.0, 0.9, size=11, color=RGBColor(0xCC,0,0))
box(s3, 4.6, 4.3, 4.1, 1.5, DARK)
code_add = """# 한 줄이면 끝
graph.add_node("translate", ...)
graph.add_edge("validate", "translate")
graph.add_edge("translate", "analyze")"""
txt(s3, code_add, 4.7, 4.38, 3.9, 1.35, size=10, color=RGBColor(0xA8,0xE6,0xFF))

# 장점 3
box(s3, 9.0, 2.15, 3.9, 4.8, RGBColor(0xFF,0xF7,0xE6))
box(s3, 9.0, 2.15, 3.9, 0.5, ORANGE)
txt(s3, "③ 복잡해질 때 빛남", 9.15, 2.18, 3.6, 0.42, size=13, bold=True, color=WHITE)
txt(s3, "지금 9주차:", 9.15, 2.78, 3.6, 0.35, size=11, bold=True, color=DARK)
txt(s3, "노드 5개 → 차이 작음", 9.15, 3.15, 3.6, 0.38, size=11, color=DARK)
txt(s3, "10주차 Self-Correction:", 9.15, 3.65, 3.6, 0.38, size=11, bold=True, color=ORANGE)
txt(s3, "조건 분기 추가\n→ 노드만 연결", 9.15, 4.05, 3.6, 0.6, size=11, color=DARK)
txt(s3, "11주차 Multi-Agent:", 9.15, 4.75, 3.6, 0.38, size=11, bold=True, color=ORANGE)
txt(s3, "노드 15개 이상\n→ while이면 스파게티\nLangGraph는 그래프만 추가", 9.15, 5.15, 3.6, 0.9, size=11, color=DARK)

txt(s3, "결론: 지금은 차이가 작아 보여도 → 복잡해질수록 LangGraph가 압도적으로 유리",
    0.4, 7.1, 12.5, 0.35, size=13, bold=True, color=BLUE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════
# 슬라이드 4 — LangGraph 핵심 개념
# ══════════════════════════════════════════════
s4_concept = prs.slides.add_slide(blank)
set_bg(s4_concept, LIGHT_BG)
box(s4_concept, 0, 0, 13.33, 1.3, BLUE)
txt(s4_concept, "LangGraph 핵심 개념: Node / Edge / State", 0.5, 0.2, 12, 1.0, size=28, bold=True, color=WHITE)

concepts = [
    ("🔵 Node", "각 작업 단계\n= Python 함수 1개",
     "collect_node()\nvalidate_node()\nanalyze_node()", RGBColor(0xE8,0xF0,0xFF), BLUE),
    ("→ Edge", "실행 순서 연결\n= 다음에 어디로 갈지",
     "add_edge(\n  'analyze',\n  'compare'\n)", RGBColor(0xF0,0xFF,0xF4), GREEN),
    ("⚡ Conditional\n   Edge", "조건 분기\n= if/else를 그래프로",
     "add_conditional_edges(\n  'validate',\n  should_retry,\n  {retry/analyze}\n)", RGBColor(0xFF,0xF7,0xE6), ORANGE),
    ("📦 State", "노드 간 공유 데이터\n= 실행 중 정보 저장",
     "class TrendState:\n  repos: list\n  analysis: str\n  report: dict", RGBColor(0xF5,0xE8,0xFF), PURPLE),
]
for i, (title, desc, code, bg, tc) in enumerate(concepts):
    x = 0.35 + i * 3.2
    box(s4_concept, x, 1.5, 3.0, 5.5, bg)
    txt(s4_concept, title, x+0.12, 1.65, 2.76, 0.7, size=15, bold=True, color=tc)
    txt(s4_concept, desc,  x+0.12, 2.4,  2.76, 0.8, size=12, color=DARK)
    box(s4_concept, x+0.1, 3.3, 2.8, 3.4, DARK)
    txt(s4_concept, code,  x+0.2, 3.4, 2.6, 3.2, size=10, color=RGBColor(0xA8,0xE6,0xFF))


# ══════════════════════════════════════════════
# 슬라이드 5 — 에이전트 실행 순서도 설계
# ══════════════════════════════════════════════
s5_flow = prs.slides.add_slide(blank)
set_bg(s5_flow, LIGHT_BG)
box(s5_flow, 0, 0, 13.33, 1.3, BLUE)
txt(s5_flow, "에이전트 실행 순서도 설계", 0.5, 0.2, 12, 1.0, size=30, bold=True, color=WHITE)

nodes = [
    ("collect",  "GitHub API\n데이터 수집",    RGBColor(0xE8,0xF0,0xFF), BLUE,   1.5),
    ("validate", "데이터\n충분한지 검증",       RGBColor(0xF0,0xFF,0xF4), GREEN,  2.7),
    ("analyze",  "AI 기술\n트렌드 분석",        RGBColor(0xFF,0xF7,0xE6), ORANGE, 3.9),
    ("compare",  "이전 기록과\n비교",           RGBColor(0xF5,0xE8,0xFF), PURPLE, 5.1),
    ("report",   "리포트 생성\n& 저장",         RGBColor(0xE5,0xFF,0xED), GREEN,  6.3),
]
for i, (name, desc, bg, tc, y) in enumerate(nodes):
    box(s5_flow, 0.5, y, 3.5, 0.9, bg)
    txt(s5_flow, name, 0.65, y+0.08, 1.5, 0.4, size=14, bold=True, color=tc)
    txt(s5_flow, desc, 2.2,  y+0.1,  1.7, 0.7, size=11, color=DARK)
    if i < 4:
        txt(s5_flow, "↓", 1.6, y+0.88, 0.5, 0.3, size=16, bold=True, color=BLUE)

box(s5_flow, 4.3, 2.55, 5.2, 1.5, RGBColor(0xFF,0xF0,0xE0))
txt(s5_flow, "조건 분기 (conditional_edge)", 4.45, 2.65, 4.9, 0.4, size=13, bold=True, color=ORANGE)
txt(s5_flow, "데이터 < 5개  →  collect로 되돌아감 (최대 3회)", 4.45, 3.1, 4.9, 0.4, size=12, color=DARK)
txt(s5_flow, "데이터 ≥ 5개  →  analyze로 진행", 4.45, 3.5, 4.9, 0.4, size=12, color=DARK)
txt(s5_flow, "↑ 이게 LangGraph의 핵심!", 4.45, 3.85, 4.9, 0.35, size=11, bold=True, color=ORANGE)

desc_items = [
    ("노드 5개", "각각 독립적인 함수로 분리"),
    ("엣지 4개", "순서 연결 (add_edge)"),
    ("조건 엣지", "validate → retry or analyze"),
    ("State",    "repos, analysis, report 공유"),
]
txt(s5_flow, "구성 요소", 9.8, 1.5, 3.3, 0.4, size=14, bold=True, color=DARK)
for i, (k, v) in enumerate(desc_items):
    y = 2.0 + i * 1.1
    box(s5_flow, 9.8, y, 3.3, 0.9, WHITE)
    box(s5_flow, 9.8, y, 1.3, 0.9, BLUE)
    txt(s5_flow, k, 9.88, y+0.2, 1.14, 0.5, size=11, bold=True, color=WHITE)
    txt(s5_flow, v, 11.2, y+0.15, 1.8, 0.6, size=10, color=DARK)


# ══════════════════════════════════════════════
# 슬라이드 6 — Graph Visualizing
# ══════════════════════════════════════════════
s6_viz = prs.slides.add_slide(blank)
set_bg(s6_viz, LIGHT_BG)
box(s6_viz, 0, 0, 13.33, 1.3, BLUE)
txt(s6_viz, "Graph Visualizing: 그래프 구조 시각화", 0.5, 0.2, 12, 1.0, size=28, bold=True, color=WHITE)

box(s6_viz, 0.4, 1.45, 5.8, 5.6, DARK)
txt(s6_viz, "ASCII 아트 (앱에서 바로 확인)", 0.55, 1.55, 5.5, 0.4, size=12, bold=True, color=GRAY)
ascii_art = """┌─────────────┐
│   collect   │ ← GitHub 수집
└──────┬──────┘
       ↓
┌──────┴──────┐
│  validate   │ ← 검증
└──────┬──────┘
  부족 │     ↘ (재수집)
       │   [collect]
  충분 ↓
┌──────┴──────┐
│   analyze   │ ← AI 분석
└──────┬──────┘
       ↓
┌──────┴──────┐
│   compare   │ ← 비교
└──────┬──────┘
       ↓
┌──────┴──────┐
│   report    │ ← 저장
└─────────────┘"""
txt(s6_viz, ascii_art, 0.55, 2.05, 5.6, 4.8, size=11, color=RGBColor(0x58,0xA6,0xFF))

box(s6_viz, 6.8, 1.45, 6.1, 5.6, DARK)
txt(s6_viz, "Mermaid 다이어그램 코드", 6.95, 1.55, 5.8, 0.4, size=12, bold=True, color=GRAY)
mermaid = """graph TD
  __start__ --> collect
  collect --> validate
  validate -->|retry| collect
  validate -->|analyze| analyze
  analyze --> compare
  compare --> report
  report --> __end__"""
txt(s6_viz, mermaid, 6.95, 2.05, 5.8, 2.5, size=12, color=RGBColor(0xA8,0xE6,0xFF))
box(s6_viz, 6.8, 4.6, 6.1, 0.5, RGBColor(0x2A,0x2A,0x40))
txt(s6_viz, "→ mermaid.live 에 붙여넣으면 예쁜 다이어그램으로 변환",
    6.95, 4.65, 5.8, 0.38, size=11, color=ORANGE)

txt(s6_viz, "LangGraph는 자기 구조를 직접 그림으로 그려준다 — 코드 안 읽어도 흐름 파악 가능",
    0.4, 7.1, 12.5, 0.35, size=13, bold=True, color=BLUE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════
# 슬라이드 7 — 만든 것: GitHub Trend Analyzer
# ══════════════════════════════════════════════
s7_impl = prs.slides.add_slide(blank)
set_bg(s7_impl, LIGHT_BG)
box(s7_impl, 0, 0, 13.33, 1.3, BLUE)
txt(s7_impl, "구현: GitHub Tech Trend Analyzer", 0.5, 0.2, 12, 1.0, size=28, bold=True, color=WHITE)

# 파일 구조
files = [
    ("github_tools.py", "GitHub Search API 호출\n언어/기간 필터 지원"),
    ("storage.py",      "분석 결과 JSON 저장\n히스토리 30개 유지"),
    ("graph.py",        "LangGraph 노드 + 엣지\n그래프 시각화 함수"),
    ("app.py",          "Streamlit 대시보드\n자동 실행 스케줄러"),
]
txt(s7_impl, "파일 구조", 0.4, 1.4, 3.0, 0.4, size=14, bold=True, color=DARK)
for i, (fname, desc) in enumerate(files):
    x = 0.4 + i * 3.2
    box(s7_impl, x, 1.85, 3.0, 1.5, WHITE)
    box(s7_impl, x, 1.85, 3.0, 0.45, BLUE)
    txt(s7_impl, fname, x+0.1, 1.88, 2.8, 0.4, size=11, bold=True, color=WHITE)
    txt(s7_impl, desc,  x+0.1, 2.4,  2.8, 0.85, size=10, color=DARK)

txt(s7_impl, "Streamlit 대시보드 기능", 0.4, 3.55, 6.0, 0.4, size=14, bold=True, color=DARK)
features = [
    "언어 필터 (Python, JS, Rust, Go 등 13개)",
    "기간 선택 (오늘 / 이번 주 / 이번 달)",
    "트렌딩 레포 목록 + 스타 수",
    "언어 분포 차트",
    "AI 트렌드 분석 + 핵심 인사이트 5개",
    "이전 분석과 비교 (변화 파악)",
    "매일 자동 실행 스케줄 설정",
    "LangGraph 그래프 구조 시각화",
]
for i, f in enumerate(features):
    col = i // 4
    row = i % 4
    x = 0.4 + col * 6.3
    y = 4.05 + row * 0.68
    box(s7_impl, x, y, 6.0, 0.55, RGBColor(0xE8,0xF0,0xFF))
    txt(s7_impl, f"✅  {f}", x+0.15, y+0.1, 5.7, 0.38, size=12, color=DARK)

box(s7_impl, 0.4, 6.85, 12.5, 0.4, BLUE)
txt(s7_impl, "실제 데이터: GitHub Search API  |  실시간 트렌딩 레포  |  AI 분석: OpenAI gpt-4o-mini",
    0.5, 6.88, 12.3, 0.35, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════
# 슬라이드 8 — 트렌드 결론 도출 과정
# ══════════════════════════════════════════════
s8_trend = prs.slides.add_slide(blank)
set_bg(s8_trend, LIGHT_BG)
box(s8_trend, 0, 0, 13.33, 1.3, BLUE)
txt(s8_trend, "트렌드 결론은 어떻게 도출되나?", 0.5, 0.2, 12, 1.0, size=30, bold=True, color=WHITE)

# STEP 1
box(s8_trend, 0.4, 1.45, 2.8, 3.8, RGBColor(0xE8,0xF0,0xFF))
box(s8_trend, 0.4, 1.45, 2.8, 0.5, BLUE)
txt(s8_trend, "STEP 1", 0.55, 1.48, 2.5, 0.42, size=13, bold=True, color=WHITE)
txt(s8_trend, "숫자 수집", 0.55, 2.05, 2.5, 0.45, size=14, bold=True, color=BLUE)
txt(s8_trend, "GitHub API", 0.55, 2.55, 2.5, 0.35, size=11, color=DARK)
rows1 = ["ollama  ⭐12,400  Go", "swarm   ⭐800  Python", "phi     ⭐500  Python"]
for i, r in enumerate(rows1):
    txt(s8_trend, r, 0.55, 2.95 + i*0.38, 2.5, 0.35, size=10, color=DARK)
txt(s8_trend, "AI 없음\n단순 목록", 0.55, 4.2, 2.5, 0.7, size=11, bold=True, color=GRAY)

txt(s8_trend, "→", 3.3, 2.9, 0.5, 0.5, size=22, bold=True, color=BLUE)

box(s8_trend, 3.9, 1.45, 2.8, 3.8, RGBColor(0xF0,0xFF,0xF4))
box(s8_trend, 3.9, 1.45, 2.8, 0.5, GREEN)
txt(s8_trend, "STEP 2", 4.05, 1.48, 2.5, 0.42, size=13, bold=True, color=WHITE)
txt(s8_trend, "통계 집계", 4.05, 2.05, 2.5, 0.45, size=14, bold=True, color=GREEN)
txt(s8_trend, "코드가 처리", 4.05, 2.55, 2.5, 0.35, size=11, color=DARK)
rows2 = ["Python  → 10개", "Go      →  5개", "Rust    →  4개", "", "ai    → 14번", "llm   →  9번"]
for i, r in enumerate(rows2):
    txt(s8_trend, r, 4.05, 2.95 + i*0.33, 2.5, 0.3, size=10, color=DARK)
txt(s8_trend, "AI 없음\n단순 집계", 4.05, 4.2, 2.5, 0.7, size=11, bold=True, color=GRAY)

txt(s8_trend, "→", 6.8, 2.9, 0.5, 0.5, size=22, bold=True, color=BLUE)

box(s8_trend, 7.4, 1.45, 2.8, 3.8, RGBColor(0xFF,0xF7,0xE6))
box(s8_trend, 7.4, 1.45, 2.8, 0.5, ORANGE)
txt(s8_trend, "STEP 3", 7.55, 1.48, 2.5, 0.42, size=13, bold=True, color=WHITE)
txt(s8_trend, "AI 해석", 7.55, 2.05, 2.5, 0.45, size=14, bold=True, color=ORANGE)
txt(s8_trend, "GPT-4o-mini", 7.55, 2.55, 2.5, 0.35, size=11, color=DARK)
txt(s8_trend, "스타 많음\n→ 관심 높음", 7.55, 2.95, 2.5, 0.6, size=10, color=DARK)
txt(s8_trend, "토픽 반복\n→ 분야 부상", 7.55, 3.6, 2.5, 0.6, size=10, color=DARK)
txt(s8_trend, "언어 분포\n→ 트렌드 파악", 7.55, 4.25, 2.5, 0.6, size=10, color=DARK)
txt(s8_trend, "AI 학습 지식\n+ 데이터 조합", 7.55, 4.75, 2.5, 0.4, size=10, bold=True, color=GRAY)

txt(s8_trend, "→", 10.3, 2.9, 0.5, 0.5, size=22, bold=True, color=BLUE)

box(s8_trend, 10.9, 1.45, 2.0, 3.8, RGBColor(0xF5,0xE8,0xFF))
box(s8_trend, 10.9, 1.45, 2.0, 0.5, PURPLE)
txt(s8_trend, "결론", 11.0, 1.48, 1.8, 0.42, size=13, bold=True, color=WHITE)
txt(s8_trend, "트렌드\n리포트", 11.0, 2.1, 1.8, 0.7, size=14, bold=True, color=PURPLE)
txt(s8_trend, '"이번 주\nLLM 에이전트\n관련 프로젝트\n급증 중"', 11.0, 2.9, 1.8, 1.6, size=10, color=DARK)

box(s8_trend, 0.4, 5.45, 12.5, 1.7, RGBColor(0xFF,0xF0,0xE0))
txt(s8_trend, "⚠️  한계 (알고 쓰기)", 0.6, 5.55, 4.0, 0.4, size=13, bold=True, color=ORANGE)
limits = [
    ("✅ 잘함", "스타 많은 레포 설명  |  언어·토픽 패턴 해석  |  읽기 좋은 문장 생성"),
    ("❌ 못함", "진짜 좋은 프로젝트인지 판단  |  스타 조작 감지  |  장기 트렌드 예측"),
]
for i, (k, v) in enumerate(limits):
    y = 6.05 + i * 0.5
    txt(s8_trend, k, 0.6, y, 1.2, 0.42, size=12, bold=True,
        color=GREEN if k.startswith("✅") else RGBColor(0xCC,0,0))
    txt(s8_trend, v, 1.9, y, 10.8, 0.42, size=12, color=DARK)

txt(s8_trend, "결론: AI가 트렌드를 발견하는 게 아니라 → 숫자 데이터를 사람이 읽기 좋은 언어로 바꿔주는 것",
    0.4, 7.1, 12.5, 0.35, size=13, bold=True, color=BLUE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════
# 슬라이드 9 — 회고 & 다음 주
# ══════════════════════════════════════════════
s9 = prs.slides.add_slide(blank)
set_bg(s9, DARK)
box(s9, 0, 0, 13.33, 1.3, BLUE)
txt(s9, "회고 & 다음 주 (10주차)", 0.5, 0.2, 12, 1.0, size=30, bold=True, color=WHITE)

box(s9, 0.5, 1.5, 5.9, 4.8, RGBColor(0x2A,0x2A,0x40))
txt(s9, "✅  이번 주 배운 것", 0.7, 1.65, 5.5, 0.5, size=16, bold=True, color=BLUE)
learns = [
    "LangGraph = while 루프를 그래프로 바꾼 것",
    "Node(함수) + Edge(연결) + State(공유 데이터)",
    "conditional_edge로 조건 분기 한 줄 처리",
    "그래프 구조 자동 시각화 (ASCII / Mermaid)",
    "실제 GitHub API로 실시간 데이터 수집",
]
for i, t in enumerate(learns):
    txt(s9, f"·  {t}", 0.7, 2.35 + i*0.72, 5.5, 0.6, size=13, color=WHITE)

box(s9, 7.0, 1.5, 5.9, 4.8, RGBColor(0x2A,0x2A,0x40))
txt(s9, "🚀  다음 주 (10주차)", 7.2, 1.65, 5.5, 0.5, size=16, bold=True, color=BLUE)
txt(s9, "Self-Correction", 7.2, 2.25, 5.5, 0.55, size=20, bold=True, color=WHITE)
txt(s9, "결과 검토 및 수정 로직 (Reflection)", 7.2, 2.85, 5.5, 0.5, size=14, color=GRAY)
nexts = [
    "틀린 결과 시 재검색/수정하는 루프",
    "지금 validate 노드의 확장 개념",
    "AI가 스스로 품질을 검토하고 개선",
]
for i, t in enumerate(nexts):
    txt(s9, f"·  {t}", 7.2, 3.5 + i*0.72, 5.5, 0.6, size=13, color=GRAY)


out = r"c:\윤주원\ai study\ai_study\week09-langgraph\juwon\week09_발표.pptx"
prs.save(out)
print(f"저장 완료: {out}")
