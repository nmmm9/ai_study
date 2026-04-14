"""
7주차 발표용 PPT 생성 스크립트
실행: python make_ppt.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# ── 색상 정의 ──────────────────────────────────
BLUE      = RGBColor(0x26, 0x5D, 0xC8)   # 포인트 블루
DARK      = RGBColor(0x1E, 0x1E, 0x2E)   # 배경 다크
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GRAY      = RGBColor(0x94, 0x94, 0xB0)
LIGHT_BG  = RGBColor(0xF4, 0xF6, 0xFF)   # 슬라이드 배경


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, text, left, top, width, height,
                size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def add_box(slide, left, top, width, height, fill_color, radius=False):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]  # 완전 빈 레이아웃


# ══════════════════════════════════════════════
# 슬라이드 1 — 표지
# ══════════════════════════════════════════════
s1 = prs.slides.add_slide(blank)
set_bg(s1, DARK)

# 포인트 라인
add_box(s1, 0, 3.2, 0.12, 1.4, BLUE)

# 주차 뱃지
add_box(s1, 0.4, 1.6, 1.5, 0.45, BLUE)
add_textbox(s1, "WEEK 07", 0.4, 1.6, 1.5, 0.45,
            size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_textbox(s1, "Function Calling", 0.4, 2.15, 10, 1.0,
            size=48, bold=True, color=WHITE)
add_textbox(s1, "질문 의도에 따른 함수 호출 및 JSON 스키마 정의",
            0.4, 3.25, 10, 0.6, size=20, color=GRAY)
add_textbox(s1, "여행 플래너 AI Agent", 0.4, 4.0, 6, 0.5,
            size=16, color=BLUE, bold=True)
add_textbox(s1, "juwon  ·  2026.04.07", 0.4, 6.6, 6, 0.4,
            size=13, color=GRAY)


# ══════════════════════════════════════════════
# 슬라이드 2 — Function Calling이란?
# ══════════════════════════════════════════════
s2 = prs.slides.add_slide(blank)
set_bg(s2, LIGHT_BG)

add_box(s2, 0, 0, 13.33, 1.3, BLUE)
add_textbox(s2, "Function Calling이란?", 0.5, 0.2, 12, 1.0,
            size=30, bold=True, color=WHITE)

add_textbox(s2, "AI는 텍스트만 주고받을 수 있어서, 실시간 데이터를 스스로 가져올 수 없습니다.",
            0.6, 1.5, 12, 0.5, size=16, color=DARK)

# 비교 박스
add_box(s2, 0.6, 2.1, 5.5, 1.8, RGBColor(0xFF, 0xE5, 0xE5))
add_textbox(s2, "❌  일반 AI", 0.8, 2.2, 5.0, 0.4, size=14, bold=True, color=RGBColor(0xCC,0,0))
add_textbox(s2, '"제주도 날씨 알려줘"\n→ "저는 실시간 정보를 모릅니다"',
            0.8, 2.6, 5.0, 1.1, size=13, color=DARK)

add_box(s2, 6.8, 2.1, 5.9, 1.8, RGBColor(0xE5, 0xFF, 0xED))
add_textbox(s2, "✅  Function Calling AI", 7.0, 2.2, 5.5, 0.4,
            size=14, bold=True, color=RGBColor(0,0x99,0x44))
add_textbox(s2, '"제주도 날씨 알려줘"\n→ get_weather("제주") 호출 → 실제 날씨 반환',
            7.0, 2.6, 5.5, 1.1, size=13, color=DARK)

add_textbox(s2, "핵심:  AI는 판단만,  실제 실행은 Python 함수가 담당",
            0.6, 4.2, 12, 0.6, size=18, bold=True, color=BLUE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════
# 슬라이드 3 — 프로젝트 소개
# ══════════════════════════════════════════════
s3 = prs.slides.add_slide(blank)
set_bg(s3, LIGHT_BG)

add_box(s3, 0, 0, 13.33, 1.3, BLUE)
add_textbox(s3, "프로젝트: 여행 플래너 AI Agent", 0.5, 0.2, 12, 1.0,
            size=30, bold=True, color=WHITE)

add_textbox(s3, "사용자가 여행 관련 질문을 하면 AI가 적절한 함수를 선택해 답변합니다.",
            0.6, 1.5, 12, 0.5, size=16, color=DARK)

cards = [
    ("🌤", "날씨 조회", "Open-Meteo API\n실시간 날씨 데이터"),
    ("🗺", "관광지 검색", "도시별 명소/맛집\n카테고리 필터"),
    ("💰", "예산 계산", "숙박/식비/교통\n항목별 상세 계산"),
    ("📅", "여행 시기", "계절별 특징\n최적 방문 시기"),
]
for i, (icon, title, desc) in enumerate(cards):
    x = 0.5 + i * 3.2
    add_box(s3, x, 2.2, 2.9, 3.2, WHITE)
    add_textbox(s3, icon,   x+0.1, 2.3,  2.7, 0.7, size=30, align=PP_ALIGN.CENTER)
    add_textbox(s3, title,  x+0.1, 3.1,  2.7, 0.5, size=15, bold=True,
                color=BLUE, align=PP_ALIGN.CENTER)
    add_textbox(s3, desc,   x+0.1, 3.65, 2.7, 1.5, size=12,
                color=GRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════
# 슬라이드 4 — 전체 구조
# ══════════════════════════════════════════════
s4 = prs.slides.add_slide(blank)
set_bg(s4, LIGHT_BG)

add_box(s4, 0, 0, 13.33, 1.3, BLUE)
add_textbox(s4, "전체 구조 (Architecture)", 0.5, 0.2, 12, 1.0,
            size=30, bold=True, color=WHITE)

steps = [
    ("사용자 질문", "main.py", "나: 제주도 날씨 어때?"),
    ("AI 판단", "agent.py", "get_weather 선택\n인자: {city: '제주'}"),
    ("함수 실행", "tools.py", "Open-Meteo API 호출\n실시간 날씨 수집"),
    ("최종 답변", "agent.py", "AI가 결과를\n자연어로 정리"),
]
colors = [
    RGBColor(0xE8,0xF0,0xFF),
    RGBColor(0xD4,0xE8,0xFF),
    RGBColor(0xC0,0xE0,0xFF),
    RGBColor(0xA8,0xD4,0xFF),
]
for i, (title, file, desc) in enumerate(steps):
    x = 0.4 + i * 3.15
    add_box(s4, x, 1.6, 2.8, 3.5, colors[i])
    add_textbox(s4, f"Step {i+1}", x+0.1, 1.7, 2.6, 0.4,
                size=12, color=BLUE, bold=True)
    add_textbox(s4, title, x+0.1, 2.1, 2.6, 0.5,
                size=16, bold=True, color=DARK)
    add_box(s4, x+0.15, 2.65, 2.5, 0.38, BLUE)
    add_textbox(s4, file, x+0.15, 2.65, 2.5, 0.38,
                size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s4, desc, x+0.1, 3.15, 2.6, 1.7,
                size=12, color=DARK)
    if i < 3:
        add_textbox(s4, "→", x+2.85, 2.9, 0.4, 0.5,
                    size=22, bold=True, color=BLUE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════
# 슬라이드 5 — 파일별 역할
# ══════════════════════════════════════════════
s5 = prs.slides.add_slide(blank)
set_bg(s5, LIGHT_BG)

add_box(s5, 0, 0, 13.33, 1.3, BLUE)
add_textbox(s5, "파일별 역할", 0.5, 0.2, 12, 1.0,
            size=30, bold=True, color=WHITE)

files = [
    ("main.py",        "대화창",    "터미널 입력 받기 · 답변 출력 · 대화 히스토리 유지"),
    ("agent.py",       "두뇌",      "OpenAI API 호출 · 함수 선택 확인 · 결과 전달 · 최종 답변 생성"),
    ("tools.py",       "도구 상자", "함수 4개 구현 · JSON 스키마 정의 · 함수 실행 디스패처"),
    (".env",           "비밀 키",   "OpenAI API 키 저장 · .gitignore 등록으로 GitHub 노출 방지"),
    ("requirements.txt","설치 목록","openai · python-dotenv · requests"),
]
for i, (fname, role, desc) in enumerate(files):
    y = 1.5 + i * 1.06
    add_box(s5, 0.5, y, 12.3, 0.88, WHITE)
    # 파일명 뱃지
    add_box(s5, 0.5, y, 2.6, 0.88, BLUE)
    add_textbox(s5, fname, 0.55, y+0.18, 2.5, 0.5,
                size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # 역할 태그
    add_box(s5, 3.2, y+0.22, 1.4, 0.42, RGBColor(0xE8,0xF0,0xFF))
    add_textbox(s5, role, 3.22, y+0.24, 1.38, 0.38,
                size=12, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    # 설명
    add_textbox(s5, desc, 4.75, y+0.2, 7.8, 0.5,
                size=12, color=DARK)

add_textbox(s5, "main.py → agent.py → tools.py → (역순) → 출력",
            0.5, 6.85, 12.3, 0.45, size=14, bold=True,
            color=BLUE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════
# 슬라이드 6 — JSON 스키마 핵심
# ══════════════════════════════════════════════
s6_json = prs.slides.add_slide(blank)
set_bg(s6_json, LIGHT_BG)

add_box(s6_json, 0, 0, 13.33, 1.3, BLUE)
add_textbox(s6_json, "핵심: JSON 스키마 설계", 0.5, 0.2, 12, 1.0,
            size=30, bold=True, color=WHITE)

add_textbox(s6_json, "AI는 JSON 스키마를 읽고 어떤 함수를 써야 할지 판단합니다.",
            0.6, 1.45, 12, 0.45, size=15, color=DARK)

# 코드 박스
add_box(s6_json, 0.6, 1.95, 7.5, 4.3, DARK)
code = '''{
  "name": "get_weather",
  "description": "특정 도시의 현재 날씨를 조회합니다.
    날씨가 궁금하거나 우산이 필요한지
    물어볼 때 사용합니다.",
  "parameters": {
    "city": {
      "type": "string",
      "description": "도시 이름 (예: 제주, 부산)"
    }
  }
}'''
add_textbox(s6_json, code, 0.8, 2.05, 7.1, 4.1, size=12,
            color=RGBColor(0xA8,0xE6,0xFF))

# 포인트
points = [
    ("name", "함수 이름 — AI가 호출할 때 사용"),
    ("description", "가장 중요! AI가 이 설명을 보고\n함수 선택 여부를 결정"),
    ("parameters", "AI가 자동으로 채워넣을 인자 정의"),
]
for i, (key, val) in enumerate(points):
    y = 2.1 + i * 1.3
    add_box(s6_json, 8.4, y, 4.4, 1.15, WHITE)
    add_textbox(s6_json, key, 8.55, y+0.05, 4.1, 0.38,
                size=13, bold=True, color=BLUE)
    add_textbox(s6_json, val, 8.55, y+0.42, 4.1, 0.65,
                size=12, color=DARK)


# ══════════════════════════════════════════════
# 슬라이드 6 — 실행 방법
# ══════════════════════════════════════════════
s6 = prs.slides.add_slide(blank)
set_bg(s6, LIGHT_BG)

add_box(s6, 0, 0, 13.33, 1.3, BLUE)
add_textbox(s6, "실행 방법", 0.5, 0.2, 12, 1.0,
            size=30, bold=True, color=WHITE)

steps_run = [
    ("1", "라이브러리 설치", "pip install -r requirements.txt"),
    ("2", "API 키 입력", ".env 파일에 OpenAI API 키 입력\nOPENAI_API_KEY=sk-..."),
    ("3", "실행", "python main.py"),
]
for i, (num, title, cmd) in enumerate(steps_run):
    y = 1.6 + i * 1.7
    add_box(s6, 0.5, y, 0.55, 0.55, BLUE)
    add_textbox(s6, num, 0.5, y, 0.55, 0.55,
                size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s6, title, 1.2, y+0.05, 3.5, 0.45,
                size=16, bold=True, color=DARK)
    add_box(s6, 1.2, y+0.6, 11.6, 0.85, DARK)
    add_textbox(s6, cmd, 1.4, y+0.65, 11.2, 0.75,
                size=13, color=RGBColor(0xA8, 0xE6, 0xFF))

add_textbox(s6, "지원 질문 예시", 0.5, 6.55, 3.0, 0.4,
            size=13, bold=True, color=BLUE)
examples = ["제주도 날씨 어때?", "부산 관광지 추천해줘", "3박 4일 예산 얼마야?", "서울 언제 가면 좋아?"]
for i, ex in enumerate(examples):
    add_box(s6, 0.5 + i*3.2, 6.95, 3.0, 0.38, RGBColor(0xE8,0xF0,0xFF))
    add_textbox(s6, ex, 0.55 + i*3.2, 6.98, 2.9, 0.32,
                size=12, color=DARK, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════
# 슬라이드 7 — 실행 예시 (데모)
# ══════════════════════════════════════════════
s7 = prs.slides.add_slide(blank)
set_bg(s7, LIGHT_BG)

add_box(s7, 0, 0, 13.33, 1.3, BLUE)
add_textbox(s7, "실행 예시 (Demo)", 0.5, 0.2, 12, 1.0,
            size=30, bold=True, color=WHITE)

# 터미널 창
add_box(s7, 0.5, 1.4, 12.3, 5.7, DARK)

# 터미널 상단 바
add_box(s7, 0.5, 1.4, 12.3, 0.4, RGBColor(0x3A,0x3A,0x5A))
add_textbox(s7, "● ● ●", 0.7, 1.42, 2.0, 0.35, size=11, color=GRAY)
add_textbox(s7, "터미널", 5.5, 1.42, 3.0, 0.35, size=11,
            color=GRAY, align=PP_ALIGN.CENTER)

terminal = [
    ("GRAY",  "나: 3박 4일 예산 얼마야?"),
    ("GRAY",  ""),
    ("GRAY",  "AI: 여행 예산을 계산하기 위해 숙박 유형, 식비 수준,"),
    ("GRAY",  "    이동 수단을 알려주시면 계산해드립니다."),
    ("GRAY",  ""),
    ("GRAY",  "나: 숙박은 저렴, 식비는 여유, 이동은 대중교통으로 해줘"),
    ("GRAY",  ""),
    ("BLUE",  "  ┌─ [함수 호출] calculate_budget"),
    ("BLUE",  "  └─ [인자]     {days: 4, accommodation: '저렴', meal: '여유', transport: '대중교통'}"),
    ("GRAY",  ""),
    ("WHITE", "AI: 3박 4일 여행 예산은 다음과 같습니다:"),
    ("WHITE", "    · 숙박비:      120,000원  (저렴 · 3박)"),
    ("WHITE", "    · 식비:        320,000원  (여유 · 4일)"),
    ("WHITE", "    · 교통비:       60,000원  (대중교통 · 4일)"),
    ("WHITE", "    · 활동/입장료:  80,000원  (4일)"),
    ("WHITE", "    ──────────────────────────"),
    ("WHITE", "    총 예산: 580,000원  |  1일 평균: 145,000원/일"),
]
color_map = {
    "GRAY": GRAY,
    "BLUE": RGBColor(0x64, 0xB5, 0xFF),
    "WHITE": WHITE,
}
for i, (clr, line) in enumerate(terminal):
    add_textbox(s7, line, 0.75, 1.92 + i*0.22, 11.8, 0.22,
                size=11, color=color_map[clr])


# ══════════════════════════════════════════════
# 슬라이드 8 — 배운 점 & 다음 주
# ══════════════════════════════════════════════
s8 = prs.slides.add_slide(blank)
set_bg(s8, DARK)

add_box(s8, 0, 0, 13.33, 1.3, BLUE)
add_textbox(s8, "회고 & 다음 주", 0.5, 0.2, 12, 1.0,
            size=30, bold=True, color=WHITE)

# 배운 점
add_box(s8, 0.5, 1.5, 5.9, 4.5, RGBColor(0x2A,0x2A,0x40))
add_textbox(s8, "✅  이번 주 배운 것", 0.7, 1.6, 5.5, 0.5,
            size=16, bold=True, color=BLUE)
learns = [
    "Function Calling 전체 흐름 이해",
    "JSON 스키마 직접 설계",
    "AI는 판단, Python은 실행 — 역할 분리",
    "description 품질이 정확도를 결정",
]
for i, t in enumerate(learns):
    add_textbox(s8, f"·  {t}", 0.7, 2.2 + i*0.75, 5.5, 0.6,
                size=14, color=WHITE)

# 다음 주
add_box(s8, 7.0, 1.5, 5.8, 4.5, RGBColor(0x2A,0x2A,0x40))
add_textbox(s8, "🚀  다음 주 (8주차)", 7.2, 1.6, 5.4, 0.5,
            size=16, bold=True, color=BLUE)
add_textbox(s8, "ReAct / Plan-and-Execute", 7.2, 2.2, 5.4, 0.5,
            size=15, bold=True, color=WHITE)
nexts = [
    "단순 1회 호출 → 다단계 추론",
    "AI가 스스로 계획 세우고",
    "여러 함수를 순서대로 호출",
    "여행 플래너를 더 똑똑하게!",
]
for i, t in enumerate(nexts):
    add_textbox(s8, f"·  {t}", 7.2, 2.8 + i*0.72, 5.4, 0.6,
                size=14, color=GRAY)

out_path = r"c:\윤주원\ai study\ai_study\week07-function-calling\juwon\week07_발표.pptx"
prs.save(out_path)
print(f"저장 완료: {out_path}")
