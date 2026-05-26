"""
make_ppt.py - Week 10 Self-Correction 발표용 PPT 생성

슬라이드 구성:
1. 표지
2. 10주차에서 만든 것
3. 트렌드를 어떻게 수집하나?
4. AI 혼자 판단 못하는 문제
5. Self-Correction이란?
6. 품질 검사 기준
7. 별로면 어떻게 다시 써?
8. 메일 발송 + GitHub 업로드
9. 파일별 역할
10. 회고
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

BG     = RGBColor(0x0d, 0x11, 0x17)
CARD   = RGBColor(0x16, 0x1b, 0x22)
BLUE   = RGBColor(0x58, 0xa6, 0xff)
GREEN  = RGBColor(0x3f, 0xb9, 0x50)
RED    = RGBColor(0xf8, 0x51, 0x49)
YELLOW = RGBColor(0xd2, 0x9e, 0x22)
WHITE  = RGBColor(0xff, 0xff, 0xff)
GRAY   = RGBColor(0x8b, 0x94, 0x9e)


def set_bg(slide, prs):
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.zorder = 0


def add_title(slide, text, top=Inches(0.3)):
    txBox = slide.shapes.add_textbox(Inches(0.5), top, Inches(9), Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = BLUE


def add_text(slide, text, left, top, width, height,
             size=14, color=None, bold=False):
    if color is None:
        color = WHITE
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_card(slide, prs, text, left, top, width, height,
             accent=BLUE, text_size=13):
    box = slide.shapes.add_shape(1, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = CARD
    box.line.color.rgb = accent
    box.line.width = Pt(1.5)

    txBox = slide.shapes.add_textbox(
        left + Inches(0.12), top + Inches(0.1),
        width - Inches(0.24), height - Inches(0.15),
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(text_size)
    run.font.color.rgb = WHITE


prs = Presentation()
prs.slide_width  = Inches(10)
prs.slide_height = Inches(5.625)
blank = prs.slide_layouts[6]


# ── 슬라이드 1: 표지 ─────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, prs)

add_text(s, "Week 10", Inches(0.5), Inches(0.9), Inches(9), Inches(0.5),
         size=18, color=GRAY)
add_text(s, "Self-Correction", Inches(0.5), Inches(1.4), Inches(9), Inches(1),
         size=44, bold=True, color=BLUE)
add_text(s, "AI가 쓴 보고서가 별로면 다시 쓰게 하는 기능",
         Inches(0.5), Inches(2.5), Inches(9), Inches(0.6),
         size=20, color=WHITE)
add_text(s, "GitHub 트렌드 수집  →  AI 분석  →  품질 검사  →  메일 + GitHub 자동 발송",
         Inches(0.5), Inches(3.2), Inches(9), Inches(0.5),
         size=15, color=GRAY)
add_text(s, "GitHub Tech Trend Analyzer v2",
         Inches(0.5), Inches(4.5), Inches(9), Inches(0.4),
         size=13, color=GREEN)


# ── 슬라이드 2: 10주차에서 만든 것 ──────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, prs)
add_title(s, "10주차에서 만든 것")

add_text(s, "한 줄 요약: GitHub에서 요즘 뜨는 프로젝트들을 AI가 분석해서,\n결과가 별로면 다시 쓰고, 완성되면 메일이랑 GitHub에 자동으로 올려주는 프로그램",
         Inches(0.5), Inches(1.0), Inches(9), Inches(0.8),
         size=14, color=GRAY)

nodes = [
    ("수집",   "GitHub에서\n트렌딩 레포\n가져오기",    BLUE),
    ("검증",   "데이터가\n충분한지\n확인",              YELLOW),
    ("생성",   "AI가\n보고서\n작성",                   BLUE),
    ("검사",   "보고서\n품질\n채점",                   YELLOW),
    ("비교",   "저번\n분석이랑\n비교",                 BLUE),
    ("발송",   "메일+\nGitHub\n업로드",               GREEN),
    ("저장",   "결과\nJSON\n저장",                    GRAY),
]
for i, (name, desc, color) in enumerate(nodes):
    add_card(s, prs, f"{name}\n\n{desc}",
             Inches(0.2 + i * 1.37), Inches(1.9), Inches(1.25), Inches(2.5),
             accent=color, text_size=12)

add_text(s, "9주차와 차이: '검사' 노드 추가 → 별로면 '생성'으로 되돌아감 (최대 3회)",
         Inches(0.5), Inches(4.7), Inches(9), Inches(0.5),
         size=13, color=GREEN)


# ── 슬라이드 3: 트렌드를 어떻게 수집하나? ────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, prs)
add_title(s, "트렌드를 어떻게 수집하나?")

add_card(s, prs,
    "GitHub에 공식 트렌딩 API가 없다\n\n"
    "그래서 이렇게 우회함:\n"
    "\"최근 7일 안에 만들어진 레포 중에\n"
    " 스타(★)를 10개 이상 받은 것들을\n"
    " 스타 많은 순서로 정렬\"",
    Inches(0.4), Inches(1.1), Inches(4.4), Inches(3.2),
    accent=BLUE, text_size=14)

add_card(s, prs,
    "기간별 기준\n\n"
    "오늘    → 오늘 하루 안에 만들어진 것\n"
    "이번 주 → 7일 안에 만들어진 것\n"
    "이번 달 → 30일 안에 만들어진 것\n\n"
    "스타 = 관심 표시 (좋아요 같은 것)\n"
    "갓 만들어졌는데 스타 많음\n= 요즘 뜨는 것",
    Inches(5.1), Inches(1.1), Inches(4.4), Inches(3.2),
    accent=GREEN, text_size=14)

add_text(s, "트렌드 데이터는 GitHub API가 수집 → 이건 객관적인 숫자라 틀릴 수 없음",
         Inches(0.5), Inches(4.6), Inches(9), Inches(0.5),
         size=13, color=GRAY)


# ── 슬라이드 4: AI 혼자 판단 못하는 문제 ────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, prs)
add_title(s, "왜 Self-Correction이 필요한가?")

add_card(s, prs,
    "AI한테 \"트렌드 분석해줘\" 라고 하면\n\n"
    "가끔 이런 답이 나옴:\n\n"
    "\"요즘 AI 관련 프로젝트가 많습니다.\"\n\n"
    "→ 너무 짧고 구체적인 내용이 없음\n"
    "→ 근데 AI는 이게 별로인 줄 모름",
    Inches(0.4), Inches(1.1), Inches(4.4), Inches(3.5),
    accent=RED, text_size=14)

add_card(s, prs,
    "AI는 자기가 쓴 게\n별로인지 스스로 잘 모름\n\n"
    "그냥 물어보면\n항상 \"잘 썼습니다\" 라고 해버림\n\n"
    "그래서 사람이 직접\n기준을 코드로 정해놓고\n코드가 대신 검사해주는 것\n= Self-Correction",
    Inches(5.1), Inches(1.1), Inches(4.4), Inches(3.5),
    accent=YELLOW, text_size=14)


# ── 슬라이드 5: Self-Correction이란? ────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, prs)
add_title(s, "Self-Correction 동작 과정")

add_text(s, "비유: 학생이 과제를 제출하면 선생님이 채점하고, 별로면 고쳐오라고 돌려보내는 것",
         Inches(0.5), Inches(1.0), Inches(9), Inches(0.5),
         size=14, color=GRAY)

steps = [
    ("AI가 보고서 작성",          "트렌드 데이터를 보고\n분석 글을 씀",                              BLUE),
    ("코드가 품질 검사",          "체크리스트로\n자동 채점\n(0~100점)",                             YELLOW),
    ("70점 미만이면",             "\"레포 이름 더 써줘\"\n피드백을 붙여서\n다시 작성 요청 (최대 3회)", RED),
    ("70점 이상이면",             "통과!\n메일이랑 GitHub에\n자동 발송",                            GREEN),
]
for i, (title, desc, color) in enumerate(steps):
    add_card(s, prs, f"{title}\n\n{desc}",
             Inches(0.3 + i * 2.37), Inches(1.7), Inches(2.1), Inches(3.0),
             accent=color, text_size=13)

add_text(s, "트렌드가 뭔지는 GitHub API가 판단 / 보고서를 잘 썼는지는 Self-Correction이 판단",
         Inches(0.5), Inches(4.9), Inches(9), Inches(0.4),
         size=12, color=GRAY)


# ── 슬라이드 6: 품질 검사 기준 ──────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, prs)
add_title(s, "품질 검사 기준 (체크리스트)")

add_text(s, "사람이 미리 정해둔 \"좋은 보고서의 조건\"을 코드로 만든 것",
         Inches(0.5), Inches(1.0), Inches(9), Inches(0.4),
         size=14, color=GRAY)

criteria = [
    ("글자 수 300자 이상",      "너무 짧으면 내용이 없는 것",                  "20점", GREEN),
    ("트렌드 키워드 3개 이상",  "\"증가, 급부상, 인기\" 같은 말 없으면 분석 아님", "25점", GREEN),
    ("레포 이름 3개 이상 언급", "구체적인 프로젝트 없으면 추상적인 말만 한 것",   "25점", GREEN),
    ("인사이트 3개 이상",       "결론이 충분히 있어야 함",                     "15점", BLUE),
    ("기술 방향성 언급",        "\"앞으로 어떻게 될 것\" 이 없으면 분석 아님",   "15점", BLUE),
]
for i, (name, reason, pts, color) in enumerate(criteria):
    y = Inches(1.55 + i * 0.75)
    add_card(s, prs, f"{pts}  |  {name}  —  {reason}",
             Inches(0.8), y, Inches(8.4), Inches(0.62),
             accent=color, text_size=13)

add_card(s, prs, "합계 70점 이상 → 통과   |   70점 미만 → 피드백 붙여서 재작성 (최대 3회)",
         Inches(0.5), Inches(5.1), Inches(9.0), Inches(0.35),
         accent=GREEN, text_size=13)


# ── 슬라이드 7: 별로면 어떻게 다시 써? ──────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, prs)
add_title(s, "별로면 어떻게 다시 써?")

add_card(s, prs,
    "1번째 시도\n\n"
    "AI 결과: \"요즘 AI 관련 프로젝트가 많습니다.\"\n\n"
    "검사 결과:\n"
    "❌ 210자 (300자 미만)\n"
    "❌ 프로젝트 이름 1개\n"
    "→ 45점 → 탈락",
    Inches(0.3), Inches(1.1), Inches(3.0), Inches(3.8),
    accent=RED, text_size=12)

add_text(s, "→\n피드백\n전달",
         Inches(3.4), Inches(2.4), Inches(0.8), Inches(1.0),
         size=13, color=YELLOW)

add_card(s, prs,
    "2번째 시도\n\n"
    "AI한테 이렇게 요청:\n"
    "\"트렌드 분석해줘\"\n"
    "+ \"저번에 이게 문제였어:\n"
    "  글자 수 짧아, 프로젝트 이름\n"
    "  3개 이상 써줘\"\n\n"
    "→ AI가 피드백 보고 다시 씀",
    Inches(4.3), Inches(1.1), Inches(3.0), Inches(3.8),
    accent=YELLOW, text_size=12)

add_text(s, "→\n통과!",
         Inches(7.4), Inches(2.4), Inches(0.7), Inches(1.0),
         size=13, color=GREEN)

add_card(s, prs,
    "결과\n\n"
    "\"llm-tool이 3,200개\n"
    "스타로 1위...\"\n\n"
    "→ 85점 → 통과\n"
    "→ 메일+GitHub 발송",
    Inches(8.2), Inches(1.1), Inches(1.6), Inches(3.8),
    accent=GREEN, text_size=11)


# ── 슬라이드 8: 메일 + GitHub 업로드 ────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, prs)
add_title(s, "품질 통과 후 자동 발송")

add_card(s, prs,
    "Gmail 자동 발송\n\n"
    "1. Gmail 우체국(smtp.gmail.com)에 접속\n"
    "2. 앱 비밀번호로 본인 확인\n"
    "3. 분석 결과로 HTML 이메일 만듦\n"
    "   (트렌딩 레포 표 + 인사이트\n"
    "    + AI 분석 내용)\n"
    "4. yoonjuwon0618@gmail.com 으로 발송",
    Inches(0.3), Inches(1.1), Inches(4.5), Inches(3.8),
    accent=BLUE, text_size=13)

add_card(s, prs,
    "GitHub README 자동 업로드\n\n"
    "1. GitHub API에 접속\n"
    "   (GitHub 토큰으로 본인 확인)\n"
    "2. trend-reports 레포에\n"
    "   TREND_REPORT.md 파일 있어?\n"
    "   있으면 → 내용 업데이트\n"
    "   없으면 → 새로 만들기\n"
    "3. 분석 결과 내용으로 저장 완료",
    Inches(5.2), Inches(1.1), Inches(4.5), Inches(3.8),
    accent=GREEN, text_size=13)

add_text(s, "핵심: 품질 검사 통과한 보고서만 발송됨 — 별로면 발송 안 됨",
         Inches(0.5), Inches(5.1), Inches(9), Inches(0.35),
         size=13, color=YELLOW)


# ── 슬라이드 9: 파일별 역할 ─────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, prs)
add_title(s, "파일별 역할")

files = [
    ("github_tools.py", "데이터 수집",     "GitHub API에 접속해서 트렌딩 레포 가져옴",                  BLUE),
    ("storage.py",      "저장/불러오기",   "분석 결과를 history.json에 저장하고 불러옴",                 GRAY),
    ("graph.py ⭐",     "전체 흐름 제어",  "7개 노드를 순서대로 연결. Self-Correction 루프 담당 (핵심)", YELLOW),
    ("notifier.py",     "발송 담당",       "Gmail 메일 발송 + GitHub 레포에 파일 업로드",               GREEN),
    ("app.py",          "화면 담당",       "브라우저에 보이는 대시보드. 품질 점수 + 이력 표시",           BLUE),
    (".env",            "비밀번호 보관",   "API 키, Gmail 비밀번호, GitHub 토큰 저장",                  GRAY),
]
for i, (name, role, desc, color) in enumerate(files):
    y = Inches(1.05 + i * 0.73)
    add_card(s, prs, f"{name}  |  {role}  —  {desc}",
             Inches(0.5), y, Inches(9.0), Inches(0.6),
             accent=color, text_size=13)


# ── 슬라이드 10: 회고 ───────────────────────────────────────
s = prs.slides.add_slide(blank)
set_bg(s, prs)
add_title(s, "회고 & 다음 주")

add_card(s, prs,
    "이번 주 배운 것\n\n"
    "• AI는 자기 결과가 좋은지 나쁜지 스스로 판단 못함\n"
    "• 사람이 기준을 코드로 정해놓으면 자동으로 품질이 올라감\n"
    "• 트렌드 수집(GitHub API)과 보고서 품질(Self-Correction)은 다른 문제\n"
    "• Gmail SMTP: 앱 비밀번호로 코드에서 직접 메일 발송 가능\n"
    "• GitHub API: 토큰만 있으면 파일 자동 생성/업데이트 가능",
    Inches(0.3), Inches(1.1), Inches(9.4), Inches(2.8),
    accent=BLUE)

add_card(s, prs,
    "다음 주 (Week 11): Multi-Agent\n\n"
    "• 지금까지: 에이전트 1개가 혼자 다 함\n"
    "• 다음 주:  에이전트 여러 개가 역할 분담해서 협력",
    Inches(0.3), Inches(4.1), Inches(9.4), Inches(1.3),
    accent=GREEN)


# ── 저장 ─────────────────────────────────────────────────────
OUTPUT = r"c:\윤주원\ai study\ai_study\week10-self-correction\juwon\week10_발표.pptx"
prs.save(OUTPUT)
print(f"저장 완료: {OUTPUT}")
