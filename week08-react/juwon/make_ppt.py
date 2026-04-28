"""
8주차 발표용 PPT 생성
실행: python make_ppt.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BLUE     = RGBColor(0x26, 0x5D, 0xC8)
DARK     = RGBColor(0x1E, 0x1E, 0x2E)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GRAY     = RGBColor(0x94, 0x94, 0xB0)
LIGHT_BG = RGBColor(0xF4, 0xF6, 0xFF)
GREEN    = RGBColor(0x16, 0xA3, 0x4A)
ORANGE   = RGBColor(0xEA, 0x58, 0x0C)

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def txt(slide, text, left, top, width, height,
        size=16, bold=False, color=DARK, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


# ══════════════════════════════════════════════
# 슬라이드 1 — 표지
# ══════════════════════════════════════════════
s1 = prs.slides.add_slide(blank)
set_bg(s1, DARK)

box(s1, 0, 3.0, 0.12, 1.6, BLUE)
box(s1, 0.4, 1.5, 1.8, 0.45, BLUE)
txt(s1, "WEEK 08", 0.4, 1.5, 1.8, 0.45, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txt(s1, "ReAct +", 0.4, 2.1, 10, 0.9, size=52, bold=True, color=WHITE)
txt(s1, "Plan-and-Execute", 0.4, 3.0, 12, 0.9, size=52, bold=True, color=BLUE)
txt(s1, "단계적 사고로 다단계 태스크 해결 에이전트 구축", 0.4, 4.0, 10, 0.6, size=18, color=GRAY)
txt(s1, "여행 플래너 AI Agent  —  10개 도구 자율 호출", 0.4, 4.7, 8, 0.5, size=15, bold=True, color=BLUE)
txt(s1, "juwon  ·  2026.04.14", 0.4, 6.6, 6, 0.4, size=13, color=GRAY)


# ══════════════════════════════════════════════
# 슬라이드 2 — AI Agent 구성 요소 + 디자인 패턴
# ══════════════════════════════════════════════
s2 = prs.slides.add_slide(blank)
set_bg(s2, LIGHT_BG)
box(s2, 0, 0, 13.33, 1.3, BLUE)
txt(s2, "AI Agent의 구성 요소와 디자인 패턴", 0.5, 0.2, 12, 1.0, size=30, bold=True, color=WHITE)

# 핵심 공식
box(s2, 0.5, 1.4, 12.3, 0.9, DARK)
txt(s2, "Agent  =  LLM  +  도구(Tools)  +  메모리  +  실행 전략(Design Pattern)",
    0.6, 1.5, 12.1, 0.7, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# 구성 요소 4개
components = [
    ("🧠 LLM",          "추론·판단 담당\n(GPT-4o-mini)",      RGBColor(0xE8,0xF0,0xFF), BLUE),
    ("🔧 Tools",         "외부 기능 호출\n(날씨·맛집·숙소…)",  RGBColor(0xFF,0xF7,0xE6), ORANGE),
    ("💾 Memory",        "대화 히스토리\n누적·관리",            RGBColor(0xF0,0xFF,0xF4), GREEN),
    ("🗺 Design Pattern", "실행 전략 결정\n(ReAct, Plan…)",    RGBColor(0xFF,0xE5,0xF5), RGBColor(0x9B,0x59,0xB6)),
]
for i, (title, desc, bg, tc) in enumerate(components):
    x = 0.4 + i * 3.2
    box(s2, x, 2.5, 3.0, 1.8, bg)
    txt(s2, title, x+0.12, 2.62, 2.76, 0.55, size=14, bold=True, color=tc)
    txt(s2, desc,  x+0.12, 3.18, 2.76, 0.95, size=12, color=DARK)

# 디자인 패턴 종류
txt(s2, "Agent Design Pattern 종류", 0.5, 4.5, 6.0, 0.45, size=15, bold=True, color=DARK)
patterns = [
    ("ReAct",             "생각하고 행동하며 반복 → 유연한 실행",   BLUE),
    ("Plan-and-Execute",  "먼저 계획, 그 다음 실행 → 체계적",      GREEN),
    ("Tool Use",          "외부 도구 호출 → 실제 데이터 사용",      ORANGE),
    ("CoT",               "Chain of Thought → 단계별 추론",       RGBColor(0x9B,0x59,0xB6)),
]
for i, (name, desc, tc) in enumerate(patterns):
    col = i % 2
    row = i // 2
    x = 0.4 + col * 6.4
    y = 5.05 + row * 0.88
    box(s2, x, y, 6.1, 0.72, WHITE)
    box(s2, x, y, 1.8, 0.72, tc)
    txt(s2, name, x+0.1, y+0.15, 1.6, 0.45, size=12, bold=True, color=WHITE)
    txt(s2, desc, x+1.95, y+0.18, 4.0, 0.45, size=12, color=DARK)

txt(s2, "이번 주: ReAct + Plan-and-Execute 두 패턴을 조합해서 구현",
    0.4, 6.85, 12.5, 0.4, size=14, bold=True, color=BLUE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════
# 슬라이드 3 — 7주차 vs 8주차
# ══════════════════════════════════════════════
s3 = prs.slides.add_slide(blank)
set_bg(s3, LIGHT_BG)
box(s3, 0, 0, 13.33, 1.3, BLUE)
txt(s3, "7주차 vs 8주차", 0.5, 0.2, 12, 1.0, size=30, bold=True, color=WHITE)

# 7주차 박스
box(s3, 0.5, 1.5, 5.8, 4.8, RGBColor(0xFF,0xE5,0xE5))
txt(s3, "7주차  Function Calling", 0.7, 1.65, 5.4, 0.5, size=16, bold=True, color=RGBColor(0xCC,0,0))
lines7 = [
    '사용자: "제주도 날씨 알려줘"',
    "",
    "AI: get_weather 선택",
    "→ 함수 1번 호출",
    "→ 끝",
    "",
    "✔ 단순한 질문만 해결 가능",
    "✔ 함수 1개로 1개 답변",
]
for i, line in enumerate(lines7):
    txt(s3, line, 0.7, 2.25 + i*0.48, 5.4, 0.45, size=13,
        color=RGBColor(0x7F,0,0) if line.startswith("✔") else DARK,
        bold=line.startswith("✔"))

# 8주차 박스
box(s3, 7.0, 1.5, 5.8, 4.8, RGBColor(0xE5,0xFF,0xED))
txt(s3, "8주차  ReAct + Plan-and-Execute", 7.2, 1.65, 5.4, 0.5, size=16, bold=True, color=GREEN)
lines8 = [
    '사용자: "제주 3박4일 완벽하게 짜줘"',
    "",
    "AI: 계획 수립 → 날씨 → 관광지",
    "    → 맛집 → 숙소 → 교통",
    "    → 예산 → 일정표 ...",
    "",
    "✔ 복잡한 질문 자율 해결",
    "✔ 10개 도구 스스로 조합",
]
for i, line in enumerate(lines8):
    txt(s3, line, 7.2, 2.25 + i*0.48, 5.4, 0.45, size=13,
        color=GREEN if line.startswith("✔") else DARK,
        bold=line.startswith("✔"))

txt(s3, "핵심 차이:  while 루프 하나가 추가되어 AI가 스스로 반복 판단",
    0.5, 6.5, 12.3, 0.5, size=15, bold=True, color=BLUE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════
# 슬라이드 4 — Plan-and-Execute
# ══════════════════════════════════════════════
s4 = prs.slides.add_slide(blank)
set_bg(s4, LIGHT_BG)
box(s4, 0, 0, 13.33, 1.3, BLUE)
txt(s4, "Plan-and-Execute: 계획 먼저, 실행 나중", 0.5, 0.2, 12, 1.0, size=28, bold=True, color=WHITE)

# 왼쪽: 설명
txt(s4, "사용자 질문을 받으면 AI가\n먼저 전체 계획을 JSON으로 수립", 0.5, 1.5, 5.5, 0.8, size=15, color=DARK)

steps = ["날씨 및 여행 시기 확인", "관광지 탐색",
         "맛집 검색", "숙소 추천", "교통편 확인",
         "예산 계산", "축제 및 행사 확인", "꿀팁 수집", "날짜별 일정표 생성"]
for i, s in enumerate(steps):
    col = i // 5
    row = i % 5
    x = 0.5 + col * 3.2
    y = 2.4 + row * 0.72
    box(s4, x, y, 3.0, 0.55, BLUE if i == 0 else RGBColor(0xE8,0xF0,0xFF))
    txt(s4, f"{i+1}. {s}", x+0.1, y+0.08, 2.8, 0.4,
        size=12, bold=(i==0),
        color=WHITE if i == 0 else DARK)

# 오른쪽: 코드
box(s4, 7.0, 1.5, 5.9, 5.6, DARK)
code = '''def plan_phase(user_query):
    response = client.chat(
        messages=[
            system_prompt,
            user_query
        ],
        # JSON 형식으로만 응답
        response_format={
            "type": "json_object"
        }
    )
    # 실행 계획 반환
    # {"city": "제주",
    #  "days": 4,
    #  "steps": [...]}
    return json.loads(response)'''
txt(s4, code, 7.2, 1.65, 5.6, 5.3, size=11, color=RGBColor(0xA8,0xE6,0xFF))


# ══════════════════════════════════════════════
# 슬라이드 5 — ReAct
# ══════════════════════════════════════════════
s5 = prs.slides.add_slide(blank)
set_bg(s5, LIGHT_BG)
box(s5, 0, 0, 13.33, 1.3, BLUE)
txt(s5, "ReAct: Reasoning + Acting  반복 루프", 0.5, 0.2, 12, 1.0, size=28, bold=True, color=WHITE)

# 왼쪽: 루프 다이어그램
cycle = [
    ("🧠 Thought", "다음에 뭘 해야 할지 생각", RGBColor(0xE8,0xF0,0xFF), BLUE),
    ("⚡ Action",   "적절한 도구 호출",         RGBColor(0xFF,0xF7,0xE6), ORANGE),
    ("👁 Observe",  "결과 확인 후 다음 결정",    RGBColor(0xF0,0xFF,0xF4), GREEN),
]
for i, (title, desc, bg, tc) in enumerate(cycle):
    y = 1.6 + i * 1.5
    box(s5, 0.5, y, 5.8, 1.2, bg if len(bg) == 3 else RGBColor(*bg[:3]))
    txt(s5, title, 0.7, y+0.15, 2.5, 0.5, size=18, bold=True, color=tc)
    txt(s5, desc,  3.3, y+0.25, 3.0, 0.5, size=13, color=DARK)
    if i < 2:
        txt(s5, "↓", 3.3, y+1.15, 1.0, 0.4, size=20, bold=True, color=BLUE)

txt(s5, "AI가 '완료'라고 판단할 때까지 반복", 0.5, 6.2, 5.8, 0.5,
    size=13, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

# 오른쪽: 코드
box(s5, 6.8, 1.4, 6.1, 5.7, DARK)
code2 = '''def react_execute(plan, query):

    while iteration < 25:
        # Thought
        response = client.chat(
            messages=messages,
            tools=TOOLS
        )
        # Action
        if message.tool_calls:
            result = execute_tool()
            messages.append(result)
        # Observe → 다음 반복

        # 완료 판단
        else:
            break  # 루프 종료

    return collected_results'''
txt(s5, code2, 7.0, 1.55, 5.8, 5.4, size=11, color=RGBColor(0xA8,0xE6,0xFF))


# ══════════════════════════════════════════════
# 슬라이드 6 — 왜 디자인 패턴이 필요한가
# ══════════════════════════════════════════════
s6 = prs.slides.add_slide(blank)
set_bg(s6, LIGHT_BG)
box(s6, 0, 0, 13.33, 1.3, BLUE)
txt(s6, "왜 Agent 디자인 패턴이 필요한가?", 0.5, 0.2, 12, 1.0, size=30, bold=True, color=WHITE)

# 왼쪽: 패턴 없을 때
box(s6, 0.4, 1.45, 5.8, 5.5, RGBColor(0xFF,0xE5,0xE5))
txt(s6, "❌  패턴 없이 (그냥 LLM)", 0.6, 1.6, 5.4, 0.5, size=15, bold=True, color=RGBColor(0xCC,0,0))
txt(s6, "입력 → 바로 답", 0.6, 2.15, 5.4, 0.45, size=13, color=DARK)
problems = [
    ("실제 데이터 없음", "날씨·맛집 정보가 틀릴 수 있음"),
    ("중간 수정 불가",   "틀려도 다시 고칠 기회 없음"),
    ("복잡한 작업 못함", "계획→검색→비교→정리 흐름 불가"),
]
for i, (title, desc) in enumerate(problems):
    y = 2.75 + i * 1.15
    box(s6, 0.6, y, 5.4, 0.95, RGBColor(0xFF,0xCC,0xCC))
    txt(s6, title, 0.75, y+0.08, 5.1, 0.38, size=13, bold=True, color=RGBColor(0xCC,0,0))
    txt(s6, desc,  0.75, y+0.46, 5.1, 0.38, size=12, color=DARK)

# 오른쪽: 패턴 있을 때
box(s6, 7.0, 1.45, 5.9, 5.5, RGBColor(0xE5,0xFF,0xED))
txt(s6, "✅  패턴 사용 (Agent)", 7.2, 1.6, 5.5, 0.5, size=15, bold=True, color=GREEN)
txt(s6, "계획 → 실행 → 확인 → 수정 → 결과", 7.2, 2.15, 5.5, 0.45, size=13, color=DARK)
patterns_why = [
    ("ReAct",             "나눠서 단계별 해결, 중간 수정 가능"),
    ("Plan-and-Execute",  "계획 먼저 → 방향 안 잃음"),
    ("Tool Use",          "도구 써서 실제 데이터 사용"),
]
for i, (title, desc) in enumerate(patterns_why):
    y = 2.75 + i * 1.15
    box(s6, 7.2, y, 5.5, 0.95, RGBColor(0xBB,0xF7,0xD0))
    txt(s6, title, 7.35, y+0.08, 5.1, 0.38, size=13, bold=True, color=GREEN)
    txt(s6, desc,  7.35, y+0.46, 5.1, 0.38, size=12, color=DARK)

txt(s6, "디자인 패턴은 AI를 똑똑하게 만드는 게 아니라 → 똑똑함을 제대로 쓰게 만드는 것",
    0.4, 7.05, 12.5, 0.35, size=14, bold=True, color=BLUE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════
# 슬라이드 7 — 전체 흐름
# ══════════════════════════════════════════════
s7 = prs.slides.add_slide(blank)
set_bg(s7, LIGHT_BG)
box(s7, 0, 0, 13.33, 1.3, BLUE)
txt(s7, "전체 흐름 & 파일 구조", 0.5, 0.2, 12, 1.0, size=30, bold=True, color=WHITE)

phases = [
    ("Phase 1", "Plan", "plan_phase()", "계획 수립\nJSON 반환", RGBColor(0xE8,0xF0,0xFF)),
    ("Phase 2", "ReAct", "react_execute()", "10개 도구\n자율 호출", RGBColor(0xE8,0xFF,0xF0)),
    ("Phase 3", "Output", "generate_html()", "HTML\n보고서 생성", RGBColor(0xFF,0xF7,0xE6)),
]
for i, (phase, name, func, desc, bg) in enumerate(phases):
    x = 0.4 + i * 4.3
    box(s7, x, 1.5, 4.0, 3.8, bg)
    txt(s7, phase, x+0.15, 1.65, 3.7, 0.45, size=12, bold=True, color=BLUE)
    txt(s7, name,  x+0.15, 2.1,  3.7, 0.7,  size=24, bold=True, color=DARK)
    box(s7, x+0.15, 2.85, 3.7, 0.5, DARK)
    txt(s7, func,  x+0.2,  2.88, 3.6, 0.44, size=12, bold=True, color=RGBColor(0xA8,0xE6,0xFF))
    txt(s7, desc,  x+0.15, 3.45, 3.7, 0.9,  size=13, color=DARK)
    if i < 2:
        txt(s7, "→", x+4.05, 2.9, 0.5, 0.5, size=28, bold=True, color=BLUE)

# 파일 구조
files = [
    ("tools.py",          "10개 함수 + JSON 스키마 + 디스패처"),
    ("agent.py",          "plan_phase() + react_execute()"),
    ("html_generator.py", "수집된 데이터 → 아름다운 HTML"),
    ("app.py",            "Streamlit 채팅 UI + 실시간 로그"),
]
txt(s7, "파일 구조", 0.4, 5.55, 3.0, 0.4, size=14, bold=True, color=DARK)
for i, (fname, desc) in enumerate(files):
    x = 0.4 + i * 3.2
    box(s7, x, 5.95, 3.05, 1.15, WHITE)
    box(s7, x, 5.95, 3.05, 0.38, BLUE)
    txt(s7, fname, x+0.08, 5.97, 2.9, 0.35, size=11, bold=True, color=WHITE)
    txt(s7, desc,  x+0.08, 6.4,  2.9, 0.65, size=10, color=DARK)


# ══════════════════════════════════════════════
# 슬라이드 8 — 회고 & 다음 주
# ══════════════════════════════════════════════
s8 = prs.slides.add_slide(blank)
set_bg(s8, DARK)
box(s8, 0, 0, 13.33, 1.3, BLUE)
txt(s8, "회고 & 다음 주 (9주차)", 0.5, 0.2, 12, 1.0, size=30, bold=True, color=WHITE)

# 배운 점
box(s8, 0.5, 1.5, 5.9, 4.7, RGBColor(0x2A,0x2A,0x40))
txt(s8, "✅  이번 주 배운 것", 0.7, 1.65, 5.5, 0.5, size=16, bold=True, color=BLUE)
learns = [
    "Plan-and-Execute로 계획 수립",
    "ReAct 루프로 다단계 자율 실행",
    "7주차 함수 4개 → 10개로 확장",
    "Streamlit UI로 채팅 + HTML 통합",
    "while 루프 하나가 Agent를 만든다",
]
for i, t in enumerate(learns):
    txt(s8, f"·  {t}", 0.7, 2.3 + i*0.72, 5.5, 0.6, size=13, color=WHITE)

# 다음 주
box(s8, 7.0, 1.5, 5.9, 4.7, RGBColor(0x2A,0x2A,0x40))
txt(s8, "🚀  다음 주 (9주차)", 7.2, 1.65, 5.5, 0.5, size=16, bold=True, color=BLUE)
txt(s8, "LangGraph 입문", 7.2, 2.25, 5.5, 0.55, size=18, bold=True, color=WHITE)
txt(s8, "그래프 기반 워크플로우\n(Nodes, Edges)", 7.2, 2.85, 5.5, 0.8, size=14, color=GRAY)
nexts = [
    "ReAct 루프 → 그래프 구조로 시각화",
    "에이전트 실행 순서도 설계",
    "Graph Visualizing 구현",
]
for i, t in enumerate(nexts):
    txt(s8, f"·  {t}", 7.2, 3.8 + i*0.72, 5.5, 0.6, size=13, color=GRAY)

out = r"c:\윤주원\ai study\ai_study\week08-react\juwon\week08_발표.pptx"
prs.save(out)
print(f"저장 완료: {out}")
